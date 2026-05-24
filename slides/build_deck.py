"""Build the workshop slide deck.

Dark theme, 16:9, ~35 slides. Embeds PNG diagrams from ../diagrams/png
and ../diagrams/personas/png. "datasense" footer on every slide.

Output:
    real-time-patterns-workshop.pptx

Run:
    python build_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

HERE = Path(__file__).parent
ROOT = HERE.parent
DIAG = ROOT / "diagrams" / "png"
PERSONAS = ROOT / "diagrams" / "personas" / "png"
OUT = HERE / "real-time-patterns-workshop.pptx"


# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
BG_DARK       = RGBColor(0x0B, 0x0E, 0x14)
BG_PANEL      = RGBColor(0x14, 0x19, 0x23)
BG_PANEL_2    = RGBColor(0x1C, 0x22, 0x30)
BORDER        = RGBColor(0x26, 0x2D, 0x3A)
TEXT          = RGBColor(0xE5, 0xE7, 0xEB)
MUTED         = RGBColor(0xA0, 0xA8, 0xB8)
DIM           = RGBColor(0x60, 0x68, 0x78)

ACCENT        = RGBColor(0x67, 0xE8, 0xF9)   # bright cyan
POLLING       = RGBColor(0xFB, 0xBF, 0x24)   # amber
WEBHOOK       = RGBColor(0xFB, 0x71, 0x85)   # pink-red
SSE           = RGBColor(0x4A, 0xDE, 0x80)   # green
WS            = RGBColor(0x60, 0xA5, 0xFA)   # blue
LLM           = RGBColor(0xC0, 0x84, 0xFC)   # purple

WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

# Fonts (system-available across macOS/Windows; falls back gracefully)
FONT_HEAD = "Helvetica Neue"
FONT_BODY = "Helvetica Neue"
FONT_CODE = "Menlo"


# ---------------------------------------------------------------------------
# Setup
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_blank_slide():
    """Add a slide with a fully dark background and the footer."""
    layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()    # no border
    bg.shadow.inherit = False
    return slide


def add_footer(slide, slide_no):
    # datasense wordmark left
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(3.0), Inches(0.3))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = ""
    r = p.add_run(); r.text = "datasense"
    r.font.name = FONT_BODY; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = MUTED
    r2 = p.add_run(); r2.text = "  ·  Real-Time Patterns Workshop"
    r2.font.name = FONT_BODY; r2.font.size = Pt(9); r2.font.color.rgb = DIM

    # slide number right
    nx = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3))
    nf = nx.text_frame
    nf.margin_left = nf.margin_right = nf.margin_top = nf.margin_bottom = 0
    np = nf.paragraphs[0]
    np.alignment = PP_ALIGN.RIGHT
    nr = np.add_run(); nr.text = str(slide_no)
    nr.font.name = FONT_BODY; nr.font.size = Pt(10); nr.font.color.rgb = DIM


def add_title(slide, text, top=Inches(0.6), color=TEXT):
    tx = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.0), Inches(0.8))
    tf = tx.text_frame
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT_HEAD; r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = color
    return tx


def add_subtitle(slide, text, top=Inches(1.45)):
    tx = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.0), Inches(0.5))
    tf = tx.text_frame
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(16); r.font.color.rgb = MUTED
    return tx


def add_text(slide, text, *, left=Inches(0.6), top=Inches(2.2), width=Inches(12.0),
             height=Inches(4.0), size=18, color=TEXT, bold=False, mono=False,
             align=PP_ALIGN.LEFT, italic_via=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.margin_left = tf.margin_top = 0
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if not line:
            continue
        r = p.add_run(); r.text = line
        r.font.name = FONT_CODE if mono else FONT_BODY
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic_via
        r.font.color.rgb = color
    return tx


def add_bullets(slide, items, *, left=Inches(0.6), top=Inches(2.2),
                width=Inches(12.0), height=Inches(4.5), size=18):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.margin_left = tf.margin_top = 0
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        # bullet dot
        dot = p.add_run(); dot.text = "•  "
        dot.font.name = FONT_BODY; dot.font.size = Pt(size); dot.font.color.rgb = ACCENT; dot.font.bold = True
        # text
        body = p.add_run(); body.text = item
        body.font.name = FONT_BODY; body.font.size = Pt(size); body.font.color.rgb = TEXT
    return tx


def add_image(slide, path: Path, *, left=None, top=None, width=None, height=None,
              max_width=Inches(12.0), max_height=Inches(5.0)):
    """Embed an image, scaled to fit within max_width x max_height, centred horizontally."""
    if not path.exists():
        raise FileNotFoundError(f"missing image: {path}")
    # Add at default size, then scale
    pic = slide.shapes.add_picture(str(path), 0, 0)
    # Compute scale to fit
    iw, ih = pic.width, pic.height
    sw = max_width / iw
    sh = max_height / ih
    s = min(sw, sh, 1.0)   # don't upscale
    pic.width = int(iw * s)
    pic.height = int(ih * s)
    pic.left = left if left is not None else int((SLIDE_W - pic.width) / 2)
    pic.top  = top  if top  is not None else int(Inches(2.0))
    return pic


def add_caption(slide, text, *, top=None, size=14):
    if top is None:
        top = Inches(6.6)
    tx = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.0), Inches(0.5))
    tf = tx.text_frame
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size); r.font.italic = True; r.font.color.rgb = MUTED
    return tx


def add_pattern_badge(slide, label, color, *, top=Inches(0.7), left=Inches(0.6)):
    """A small colored pill labelling the slide's pattern (e.g. POLLING)."""
    width = Inches(1.6); height = Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT_HEAD; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = BG_DARK
    return shape


def section_divider(title, color):
    slide = add_blank_slide()
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))

    # Big colored bar on the left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False

    # Big title text, vertically centered
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.5), Inches(1.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.name = FONT_HEAD; r.font.size = Pt(60); r.font.bold = True; r.font.color.rgb = TEXT

    # Small "SECTION X" label above
    sx = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.5), Inches(0.4))
    sf = sx.text_frame
    sp = sf.paragraphs[0]
    sr = sp.add_run(); sr.text = "SECTION"
    sr.font.name = FONT_HEAD; sr.font.size = Pt(14); sr.font.bold = True; sr.font.color.rgb = color
    return slide


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_cover():
    slide = add_blank_slide()

    # Top-left small label
    lx = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(6.0), Inches(0.4))
    lf = lx.text_frame; lp = lf.paragraphs[0]
    lr = lp.add_run(); lr.text = "datasense workshop"
    lr.font.name = FONT_HEAD; lr.font.size = Pt(14); lr.font.bold = True; lr.font.color.rgb = ACCENT

    # Big title - two lines, 56pt fits comfortably in 12in width
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12.0), Inches(2.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Real-Time"
    r.font.name = FONT_HEAD; r.font.size = Pt(64); r.font.bold = True; r.font.color.rgb = TEXT
    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    r2 = p2.add_run(); r2.text = "Communication Patterns"
    r2.font.name = FONT_HEAD; r2.font.size = Pt(56); r2.font.bold = True; r2.font.color.rgb = TEXT

    # Subtitle - the 4 patterns, moved lower so it doesn't collide with the title
    sx = slide.shapes.add_textbox(Inches(0.6), Inches(5.65), Inches(12.0), Inches(0.6))
    sf = sx.text_frame; sp = sf.paragraphs[0]
    for i, (label, color) in enumerate([
        ("Polling", POLLING), ("Webhooks", WEBHOOK), ("SSE", SSE), ("WebSockets", WS),
    ]):
        if i > 0:
            sep = sp.add_run(); sep.text = "   ·   "
            sep.font.name = FONT_HEAD; sep.font.size = Pt(22); sep.font.color.rgb = DIM
        r = sp.add_run(); r.text = label
        r.font.name = FONT_HEAD; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = color

    # Bottom tagline
    by = slide.shapes.add_textbox(Inches(0.6), Inches(6.45), Inches(12.0), Inches(0.4))
    bf = by.text_frame; bp = bf.paragraphs[0]
    br = bp.add_run(); br.text = "Pick the right pattern. Build apps that feel live without overengineering."
    br.font.name = FONT_BODY; br.font.size = Pt(16); br.font.italic = True; br.font.color.rgb = MUTED

    add_footer(slide, 1)


def slide_agenda():
    slide = add_blank_slide()
    add_title(slide, "What we'll cover")
    add_subtitle(slide, "Four patterns, three projects, one running scenario.")

    items = [
        "The four real-time patterns and what each is for",
        "When to pick which (the decision tree)",
        "HTTP basics + JWT auth, refreshed",
        "Hands-on: 7 runnable example folders (server + client per pattern)",
        "Three full projects, including the LiveOrder capstone that composes all four",
        "Where these show up in AI apps (LLM streaming, MCP, voice agents)",
    ]
    add_bullets(slide, items, top=Inches(2.3), size=20)
    add_footer(slide, 2)


def slide_cast():
    slide = add_blank_slide()
    add_title(slide, "Meet the cast")
    add_subtitle(slide, "Every example below is told through these four. Same people, different patterns.")
    add_image(slide, PERSONAS / "p_cast.png", top=Inches(2.0), max_height=Inches(4.4))
    add_caption(slide, "Maya builds it · Raj uses it · Priya cooks it · Sam delivers it",
                top=Inches(6.6))
    add_footer(slide, 3)


def slide_http_problem():
    slide = add_blank_slide()
    add_title(slide, "The problem with plain HTTP")
    add_subtitle(slide, "The client always speaks first. The server can never push.")
    add_image(slide, DIAG / "01_http_basic.png", top=Inches(2.0), max_height=Inches(4.6))
    add_caption(slide, "Every real-time pattern in this workshop exists to work around this one limit.",
                top=Inches(6.6))
    add_footer(slide, 4)


def slide_status_codes():
    slide = add_blank_slide()
    add_title(slide, "HTTP status codes that matter for real-time")
    add_subtitle(slide, "Memorise these. You'll handle them in every pattern below.")
    rows = [
        ("200", "OK",                    "successful poll; successful webhook receipt"),
        ("201", "Created",               "after a POST that creates something"),
        ("202", "Accepted",              "'Got it, processing async' - webhook receivers"),
        ("204", "No Content",            "long-poll timeout; nothing to send"),
        ("304", "Not Modified",          "ETag/If-None-Match - smart polling savings"),
        ("401", "Unauthorized",          "missing/expired token; bad webhook signature"),
        ("404", "Not Found",             "wrong URL or missing record"),
        ("408", "Request Timeout",       "long-poll fell past server limit"),
        ("429", "Too Many Requests",     "polling too fast or rate-limited"),
        ("500", "Internal Server Error", "bug; crashed dependency"),
        ("502", "Bad Gateway",           "proxy got nothing useful from backend"),
        ("503", "Service Unavailable",   "maintenance, overload; webhook senders retry"),
    ]
    status_table(slide, rows, top=Inches(2.3))
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_stateless_table():
    slide = add_blank_slide()
    add_title(slide, "Stateless vs stateful, per pattern")
    add_subtitle(slide, "HTTP itself is stateless. Some patterns add stateful connections on top.")
    rows = [
        ("Polling (short)",  "new HTTP request per poll",                  "Stateless"),
        ("Polling (long)",   "one held HTTP request",                      "Stateless on server, in-flight on client"),
        ("Webhook receiver", "separate HTTP request per event",            "Stateless"),
        ("SSE",              "one long-lived HTTP connection",             "Stateful at the transport level"),
        ("WebSocket",        "one persistent connection, both directions", "Stateful"),
    ]
    # 3-col bigger table
    headers = ("PATTERN", "CONNECTION MODEL", "STATEFUL?")
    col_w = [Inches(3.0), Inches(5.6), Inches(3.7)]
    top = Inches(2.5); row_h = Inches(0.7)
    x = Inches(0.6)
    for i, h in enumerate(headers):
        hx = slide.shapes.add_textbox(x, top, col_w[i], Inches(0.4))
        hr = hx.text_frame.paragraphs[0].add_run(); hr.text = h
        hr.font.name = FONT_HEAD; hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = ACCENT
        x += col_w[i]
    for ri, row in enumerate(rows):
        y = top + Inches(0.5) + row_h * ri
        x = Inches(0.6)
        for ci, cell in enumerate(row):
            tx = slide.shapes.add_textbox(x, y, col_w[ci], row_h - Inches(0.1))
            tp = tx.text_frame.paragraphs[0]
            tr = tp.add_run(); tr.text = cell
            tr.font.name = FONT_BODY; tr.font.size = Pt(14)
            tr.font.color.rgb = TEXT if ci != 2 else (SSE if "Stateful" in cell and "transport" not in cell else MUTED)
            x += col_w[ci]
    add_text(slide,
             "Stateful = costs memory per client. 100K SSE/WS connections = ~10 GB. Plan accordingly.",
             top=Inches(6.2), size=14, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_four_patterns_intro():
    slide = add_blank_slide()
    add_title(slide, "The four patterns at a glance")
    add_subtitle(slide, "Mental model first. Wire format and code in the next sections.")

    # 4 columns, one per pattern
    col_w = Inches(2.9); col_h = Inches(4.0); gap = Inches(0.2); start_x = Inches(0.7)
    top = Inches(2.5)
    patterns = [
        ("POLLING",   POLLING, "Are we there yet?",        "Client repeatedly asks the server on a timer. Crude. Sometimes the right answer."),
        ("WEBHOOKS",  WEBHOOK, "Don't call us,\nwe'll call you", "External server POSTs to your URL when something happens. Server-to-server."),
        ("SSE",       SSE,     "Tune in.\nServer broadcasts.", "Server holds one HTTP connection and pushes events. Browser-native. One-way."),
        ("WEBSOCKETS",WS,      "Open a phone line",        "Persistent, bidirectional, low-overhead. Use only when you really need two-way."),
    ]
    for i, (name, color, tagline, body) in enumerate(patterns):
        x = start_x + (col_w + gap) * i
        # card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.75)
        card.shadow.inherit = False
        # colour stripe on top
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, col_w, Inches(0.12))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = color
        stripe.line.fill.background(); stripe.shadow.inherit = False
        # title
        ttx = slide.shapes.add_textbox(x, top + Inches(0.35), col_w, Inches(0.5))
        ttf = ttx.text_frame; ttf.margin_left = Inches(0.2); ttf.margin_right = Inches(0.2); ttf.margin_top = 0
        ttp = ttf.paragraphs[0]
        ttr = ttp.add_run(); ttr.text = name
        ttr.font.name = FONT_HEAD; ttr.font.size = Pt(20); ttr.font.bold = True; ttr.font.color.rgb = color
        # tagline
        tagx = slide.shapes.add_textbox(x, top + Inches(1.0), col_w, Inches(0.9))
        tagf = tagx.text_frame; tagf.margin_left = Inches(0.2); tagf.margin_right = Inches(0.2); tagf.margin_top = 0
        for j, line in enumerate(tagline.split("\n")):
            tp = tagf.paragraphs[0] if j == 0 else tagf.add_paragraph()
            tr = tp.add_run(); tr.text = line
            tr.font.name = FONT_HEAD; tr.font.size = Pt(15); tr.font.italic = True; tr.font.color.rgb = TEXT
        # body
        bodyx = slide.shapes.add_textbox(x, top + Inches(2.2), col_w, Inches(1.7))
        bodyf = bodyx.text_frame; bodyf.margin_left = Inches(0.2); bodyf.margin_right = Inches(0.2); bodyf.margin_top = 0
        bodyf.word_wrap = True
        bp = bodyf.paragraphs[0]
        br = bp.add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(13); br.font.color.rgb = MUTED
    add_footer(slide, 5)


# ----- Reusable building blocks for "deep dive" slides -----
def code_block(slide, code: str, *, left=Inches(0.6), top=Inches(2.4),
               width=Inches(12.0), height=Inches(4.0), size=12):
    """Render a code snippet in a dark panel with monospace text + soft accent border."""
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid(); panel.fill.fore_color.rgb = BG_PANEL
    panel.line.color.rgb = BORDER; panel.line.width = Pt(0.5); panel.shadow.inherit = False
    tx = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2),
                                   width - Inches(0.5), height - Inches(0.4))
    tf = tx.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.name = FONT_CODE; r.font.size = Pt(size); r.font.color.rgb = TEXT


def two_col_cards(slide, *, top, left_title, right_title, left_color, right_color,
                  left_items, right_items, item_size=13):
    """Two cards side by side, used for 'use when / don't use when' style slides."""
    col_w = Inches(5.9); col_h = Inches(4.0); gap = Inches(0.3)
    # left card
    lcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, col_w, col_h)
    lcard.fill.solid(); lcard.fill.fore_color.rgb = BG_PANEL
    lcard.line.color.rgb = left_color; lcard.line.width = Pt(1.5); lcard.shadow.inherit = False
    lh = slide.shapes.add_textbox(Inches(0.85), top + Inches(0.2), col_w - Inches(0.5), Inches(0.5))
    lhp = lh.text_frame.paragraphs[0]
    lhr = lhp.add_run(); lhr.text = left_title
    lhr.font.name = FONT_HEAD; lhr.font.size = Pt(18); lhr.font.bold = True; lhr.font.color.rgb = left_color
    add_text(slide, "\n".join(left_items),
             left=Inches(0.85), top=top + Inches(0.85),
             width=col_w - Inches(0.5), height=Inches(3.0),
             size=item_size, color=TEXT)
    # right card
    rx = Inches(0.6) + col_w + gap
    rcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, top, col_w, col_h)
    rcard.fill.solid(); rcard.fill.fore_color.rgb = BG_PANEL
    rcard.line.color.rgb = right_color; rcard.line.width = Pt(1.5); rcard.shadow.inherit = False
    rh = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.2), col_w - Inches(0.5), Inches(0.5))
    rhp = rh.text_frame.paragraphs[0]
    rhr = rhp.add_run(); rhr.text = right_title
    rhr.font.name = FONT_HEAD; rhr.font.size = Pt(18); rhr.font.bold = True; rhr.font.color.rgb = right_color
    add_text(slide, "\n".join(right_items),
             left=rx + Inches(0.25), top=top + Inches(0.85),
             width=col_w - Inches(0.5), height=Inches(3.0),
             size=item_size, color=TEXT)


def status_table(slide, rows, *, top=Inches(2.4), col_widths=None):
    """A clean status-code-style table on a dark slide. rows = list of (code, name, where)."""
    if col_widths is None:
        col_widths = [Inches(1.3), Inches(2.6), Inches(8.1)]
    headers = ("CODE", "MEANING", "WHEN YOU'LL SEE IT")
    row_h = Inches(0.36)
    # header row
    x = Inches(0.6)
    for i, h in enumerate(headers):
        hx = slide.shapes.add_textbox(x, top, col_widths[i], Inches(0.3))
        hp = hx.text_frame.paragraphs[0]
        hr = hp.add_run(); hr.text = h
        hr.font.name = FONT_HEAD; hr.font.size = Pt(11); hr.font.bold = True; hr.font.color.rgb = ACCENT
        x += col_widths[i]
    # data rows
    for ri, row in enumerate(rows):
        y = top + Inches(0.4) + row_h * ri
        x = Inches(0.6)
        for ci, cell in enumerate(row):
            tx = slide.shapes.add_textbox(x, y, col_widths[ci], row_h)
            tp = tx.text_frame.paragraphs[0]
            tr = tp.add_run(); tr.text = cell
            tr.font.name = FONT_CODE if ci == 0 else FONT_BODY
            tr.font.size = Pt(13)
            tr.font.color.rgb = ACCENT if ci == 0 else TEXT
            x += col_widths[ci]


# ----- Generic helpers for paired (concept | example) slides -----
def concept_slide(badge_label, badge_color, title_text, subtitle_text, image_path):
    """A 'how the pattern works' slide using a persona-FREE technical diagram."""
    slide = add_blank_slide()
    add_pattern_badge(slide, badge_label, badge_color)
    # Small "CONCEPT" sub-badge to make the pair contrast explicit
    cx = slide.shapes.add_textbox(Inches(2.3), Inches(0.72), Inches(1.5), Inches(0.3))
    cf = cx.text_frame; cp = cf.paragraphs[0]
    cr = cp.add_run(); cr.text = "CONCEPT"
    cr.font.name = FONT_HEAD; cr.font.size = Pt(11); cr.font.bold = True; cr.font.color.rgb = MUTED
    add_title(slide, title_text, top=Inches(1.15))
    add_subtitle(slide, subtitle_text, top=Inches(1.95))
    add_image(slide, image_path, top=Inches(2.55), max_height=Inches(4.0))
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def example_slide(badge_label, badge_color, title_text, subtitle_text, image_path):
    """A 'here's how it plays out for one of our personas' slide."""
    slide = add_blank_slide()
    add_pattern_badge(slide, badge_label, badge_color)
    # "EXAMPLE" sub-badge in the persona color so the eye learns the pair
    cx = slide.shapes.add_textbox(Inches(2.3), Inches(0.72), Inches(1.5), Inches(0.3))
    cf = cx.text_frame; cp = cf.paragraphs[0]
    cr = cp.add_run(); cr.text = "EXAMPLE"
    cr.font.name = FONT_HEAD; cr.font.size = Pt(11); cr.font.bold = True; cr.font.color.rgb = badge_color
    add_title(slide, title_text, top=Inches(1.15))
    add_subtitle(slide, subtitle_text, top=Inches(1.95))
    add_image(slide, image_path, top=Inches(2.55), max_height=Inches(4.0))
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_polling_cheaper():
    slide = add_blank_slide()
    add_pattern_badge(slide, "POLLING", POLLING)
    add_title(slide, "Making polling cheaper without changing patterns", top=Inches(1.15))
    add_subtitle(slide, "Six tricks that cut waste 10x. Use them ALL before you reach for SSE.",
                 top=Inches(1.95))
    items = [
        ("1. Cursor / since= parameter", "Send the last-seen id so the server returns only deltas."),
        ("2. ETag + If-None-Match",      "Server returns 304 with empty body when nothing changed. Bandwidth -> ~zero."),
        ("3. Exponential backoff",       "If five polls in a row return empty, slow down. Reset on activity."),
        ("4. Jitter",                    "Add +-10% random offset so 10K clients don't sync up and DDoS you."),
        ("5. Pause when invisible",      "Page Visibility API. No point polling for a tab no one's looking at."),
        ("6. Conditional auth caching",  "Don't hit your auth DB on every poll. JWT verify in-process is fine."),
    ]
    top = Inches(2.7); h = Inches(0.6); gap = Inches(0.1)
    for i, (head, body) in enumerate(items):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.0), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        hx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.08), Inches(4.0), Inches(0.4))
        hp = hx.text_frame.paragraphs[0]
        hr = hp.add_run(); hr.text = head
        hr.font.name = FONT_HEAD; hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = POLLING
        bx = slide.shapes.add_textbox(Inches(5.0), y + Inches(0.08), Inches(7.5), Inches(0.4))
        bp = bx.text_frame.paragraphs[0]
        br = bp.add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(13); br.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_polling_bugs():
    slide = add_blank_slide()
    add_pattern_badge(slide, "POLLING", POLLING)
    add_title(slide, "Polling bugs you'll hit, and the fix", top=Inches(1.15))
    add_subtitle(slide, "Every team rediscovers these. Save yourself the outage.",
                 top=Inches(1.95))
    bugs = [
        ("'Data never updates' even though server has new data",
         "Browser or CDN cached your GET. Add Cache-Control: no-store, or ?t=<ts> cache-buster."),
        ("Long-poll times out at 30s even though you set 60s",
         "A proxy or LB has a shorter idle timeout. Set server < LB always."),
        ("Worked in dev, exploded in prod",
         "Dev had 1 user; prod has 10K. Add per-endpoint req/sec and DB query metrics; you'll see the N+1."),
        ("Polling continues after tab closes",
         "Forgot to clearInterval on unmount/visibilitychange. Always pair start with stop."),
        ("UI shows duplicate notifications on every poll",
         "Missing cursor; you're re-fetching the whole list. Add ?since=<last_id>."),
    ]
    top = Inches(2.7); h = Inches(0.78); gap = Inches(0.08)
    for i, (sym, fix) in enumerate(bugs):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.0), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        sx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.1), Inches(11.5), Inches(0.3))
        sp = sx.text_frame.paragraphs[0]
        sr = sp.add_run(); sr.text = sym
        sr.font.name = FONT_HEAD; sr.font.size = Pt(13); sr.font.bold = True; sr.font.color.rgb = WEBHOOK
        fx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.4), Inches(11.5), Inches(0.4))
        fp = fx.text_frame.paragraphs[0]
        fr = fp.add_run(); fr.text = fix
        fr.font.name = FONT_BODY; fr.font.size = Pt(12); fr.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ----- Polling section -----
def polling_section():
    section_divider("Polling", POLLING)

    # PAIR 1: short polling
    concept_slide("POLLING", POLLING,
                  "Short polling - the protocol",
                  "Client asks on a timer. Server immediately returns whatever it has.",
                  DIAG / "02_short_polling.png")
    example_slide("POLLING", POLLING,
                  "Raj refreshing his Swiggy order",
                  "The kid in the back seat. 'Are we there yet?'  Most refreshes return the same status.",
                  PERSONAS / "p_polling_swiggy.png")

    # PAIR 2: long polling
    concept_slide("POLLING", POLLING,
                  "Long polling - the protocol",
                  "Server holds the request open until data is ready (or timeout). One held connection per client.",
                  DIAG / "03_long_polling.png")
    example_slide("POLLING", POLLING,
                  "Raj waits for an Uber driver",
                  "One held request. The server replies the moment a driver accepts. No wasted polls.",
                  PERSONAS / "p_long_polling_uber.png")
    slide_polling_code()
    slide_polling_cost_math()
    slide_polling_cheaper()
    slide_polling_bugs()
    slide_polling_when_right()

    # 10: trade-offs
    slide = add_blank_slide()
    add_pattern_badge(slide, "POLLING", POLLING)
    add_title(slide, "When polling is the right call (and when it isn't)", top=Inches(1.15))
    add_subtitle(slide, "It's not always wrong. Just usually mis-used.", top=Inches(1.95))
    # two-column: good vs bad
    col_top = Inches(2.7); col_w = Inches(5.9); gap = Inches(0.3)
    # good
    good = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), col_top, col_w, Inches(4.0))
    good.fill.solid(); good.fill.fore_color.rgb = BG_PANEL
    good.line.color.rgb = SSE; good.line.width = Pt(1.5); good.shadow.inherit = False
    gh = slide.shapes.add_textbox(Inches(0.85), col_top + Inches(0.2), col_w - Inches(0.5), Inches(0.5))
    ghp = gh.text_frame.paragraphs[0]
    ghr = ghp.add_run(); ghr.text = "Use polling when"
    ghr.font.name = FONT_HEAD; ghr.font.size = Pt(18); ghr.font.bold = True; ghr.font.color.rgb = SSE
    add_text(slide, "\n".join([
        "Updates are rare and a few seconds stale is fine",
        "You're calling someone else's API that only supports REST",
        "Total clients < 100; engineering cost matters more than waste",
        "You need to ship in an afternoon",
        "Status checks for slow background jobs (batch reports, fine-tunes)",
    ]), left=Inches(0.85), top=col_top + Inches(0.85), width=col_w - Inches(0.5), height=Inches(3.0),
        size=13, color=TEXT)
    # bad
    bad_x = Inches(0.6) + col_w + gap
    bad = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bad_x, col_top, col_w, Inches(4.0))
    bad.fill.solid(); bad.fill.fore_color.rgb = BG_PANEL
    bad.line.color.rgb = WEBHOOK; bad.line.width = Pt(1.5); bad.shadow.inherit = False
    bh = slide.shapes.add_textbox(bad_x + Inches(0.25), col_top + Inches(0.2), col_w - Inches(0.5), Inches(0.5))
    bhp = bh.text_frame.paragraphs[0]
    bhr = bhp.add_run(); bhr.text = "Don't poll when"
    bhr.font.name = FONT_HEAD; bhr.font.size = Pt(18); bhr.font.bold = True; bhr.font.color.rgb = WEBHOOK
    add_text(slide, "\n".join([
        "The UI must feel live (chat, typing, multiplayer, game)",
        "You have 10k+ mostly-idle clients (millions of empty requests)",
        "You'd be polling at 10+ Hz - just use SSE",
        "You're hitting someone else's rate-limited API quota",
        "The source already supports webhooks; just use those",
    ]), left=bad_x + Inches(0.25), top=col_top + Inches(0.85), width=col_w - Inches(0.5), height=Inches(3.0),
        size=13, color=TEXT)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_webhook_local_testing():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "Testing webhooks locally", top=Inches(1.15))
    add_subtitle(slide, "Stripe can't reach localhost:8000. Here's how to bridge.",
                 top=Inches(1.95))
    opts = [
        ("webhook.site",       "look only",
         "Visit https://webhook.site, get a unique URL, paste into Stripe. Sees every request in a dashboard. Can't run YOUR code against it but great for inspecting payloads."),
        ("ngrok / Cloudflare Tunnel / Tailscale Funnel", "tunnel",
         "ngrok http 8000 -> public URL forwarding to localhost. Combined with the ngrok web inspector (localhost:4040) you can see + replay every request."),
        ("Stripe CLI (or provider equivalent)",  "real events, no tunnel",
         "stripe listen --forward-to localhost:8000/webhooks/stripe.  Forwards real test webhook events from your Stripe account straight to localhost. Handles signature setup."),
    ]
    top = Inches(2.7); h = Inches(1.15); gap = Inches(0.1)
    for i, (name, kind, body) in enumerate(opts):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.0), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        hx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.15), Inches(7.5), Inches(0.35))
        hp = hx.text_frame.paragraphs[0]
        hr = hp.add_run(); hr.text = name
        hr.font.name = FONT_HEAD; hr.font.size = Pt(15); hr.font.bold = True; hr.font.color.rgb = WEBHOOK
        kx = slide.shapes.add_textbox(Inches(8.5), y + Inches(0.2), Inches(3.5), Inches(0.3))
        kp = kx.text_frame.paragraphs[0]; kp.alignment = PP_ALIGN.RIGHT
        kr = kp.add_run(); kr.text = kind
        kr.font.name = FONT_HEAD; kr.font.size = Pt(11); kr.font.italic = True; kr.font.color.rgb = MUTED
        bx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.55), Inches(11.3), Inches(0.55))
        bf = bx.text_frame; bf.word_wrap = True
        bp = bf.paragraphs[0]
        br = bp.add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(12); br.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_webhook_being_sender():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "If you're the SENDER (offering webhooks to customers)", top=Inches(1.15))
    add_subtitle(slide, "The other side of the table. Same rules, mirrored.",
                 top=Inches(1.95))
    items = [
        "Make registration first-class. UI + API to register URLs, multiple URLs per account, filter by event type.",
        "Sign every payload. Per-customer signing secret. Document the algorithm and rotation flow.",
        "Include a timestamp in the signed payload (defeats replay attacks).",
        "Retry with backoff. Standard schedule: 1m, 5m, 15m, 1h, 6h, 12h. Give up after ~24h, mark endpoint disabled.",
        "Provide a 'Send test event' button in the dashboard - your customers will use it constantly.",
        "Maintain a delivery log: last 100 deliveries, status codes, response times, response bodies. Invaluable for support.",
        "Build internal rate limits. A bug that fires 1M duplicate webhooks tonight = 1M customer outages.",
    ]
    add_bullets(slide, items, top=Inches(2.5), size=14)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ----- Webhook section -----
def webhook_section():
    section_divider("Webhooks", WEBHOOK)

    # PAIR 1: basic flow
    concept_slide("WEBHOOK", WEBHOOK,
                  "Webhook - the protocol",
                  "Third-party server POSTs to YOUR URL when something happens. Zero idle traffic.",
                  DIAG / "05_webhook_basic.png")
    example_slide("WEBHOOK", WEBHOOK,
                  "Stripe pays for Raj's biryani",
                  "Stripe is the source of truth. They tell Maya's backend when the card cleared.",
                  PERSONAS / "p_webhook_stripe.png")

    slide_webhook_role_flip()
    slide_webhook_payload()

    # the 4 rules
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "The four rules of webhook receivers", top=Inches(1.15))
    add_subtitle(slide, "Skip any one and you'll have very interesting outages.", top=Inches(1.95))
    rules = [
        ("1. Return 2xx fast", "Don't do work in the handler. Verify, queue, return 200. Stripe times out at 5-10 seconds."),
        ("2. Verify the signature", "HMAC-SHA256 of the body, constant-time compare. Without this, anyone can fake events."),
        ("3. Dedup by event id", "Senders sometimes deliver twice. Without dedup, you'll double-charge cards."),
        ("4. Return 200 (not 5xx) for events you don't care about", "5xx triggers retries forever. Log + 200 is the right answer for noise."),
    ]
    top = Inches(2.7); h = Inches(1.0); gap = Inches(0.1)
    for i, (head, body) in enumerate(rules):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.0), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        hx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.15), Inches(11.5), Inches(0.4))
        hf = hx.text_frame; hp = hf.paragraphs[0]
        hr = hp.add_run(); hr.text = head
        hr.font.name = FONT_HEAD; hr.font.size = Pt(15); hr.font.bold = True; hr.font.color.rgb = WEBHOOK
        bx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.5), Inches(11.5), Inches(0.5))
        bf = bx.text_frame; bf.word_wrap = True
        bp = bf.paragraphs[0]
        br = bp.add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(13); br.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))

    # retries + dedup
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "Retries and idempotent dedup", top=Inches(1.15))
    add_subtitle(slide, "Same event id arriving twice is normal. Be ready.", top=Inches(1.95))
    add_image(slide, DIAG / "06_webhook_retries.png", top=Inches(2.55), max_height=Inches(4.0))
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))

    # HMAC security
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "Signature verification - HMAC-SHA256", top=Inches(1.15))
    add_subtitle(slide, "The one thing standing between your webhook URL and fake events.", top=Inches(1.95))
    add_image(slide, DIAG / "07_webhook_security.png", top=Inches(2.55), max_height=Inches(4.0))
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))

    # Deep dives
    slide_webhook_receiver_code()
    slide_webhook_local_testing()
    slide_webhook_being_sender()
    slide_webhook_real_world()


def slide_sse_grammar():
    slide = add_blank_slide()
    add_pattern_badge(slide, "SSE", SSE)
    add_title(slide, "The SSE grammar - all of it", top=Inches(1.15))
    add_subtitle(slide, "Five fields plus a blank line. That's the protocol in full.",
                 top=Inches(1.95))
    code_block(slide, """\
event: <event-type>      optional, default 'message'.  JS addEventListener key.
id: <event-id>           optional. Browser remembers; sends as Last-Event-ID on reconnect.
retry: <ms>              optional. Tells browser the reconnect delay.
data: <line of payload>  can repeat; multiple data: lines join with '\\n'.
:<comment>               ignored by client. Use to send keep-alive pings.
                         <- BLANK LINE ends this event

Examples
--------
  data: hello                              # simplest possible event

  event: token                             # named event with JSON payload
  data: {"text": "Vada"}

  id: 42                                   # with id - enables Last-Event-ID resume
  event: chunk
  data: {"text": " Pav"}

  : keep-alive ping                        # comment - keeps proxies from closing
""", top=Inches(2.5), height=Inches(4.6), size=12)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_sse_gotchas():
    slide = add_blank_slide()
    add_pattern_badge(slide, "SSE", SSE)
    add_title(slide, "SSE gotchas every team hits", top=Inches(1.15))
    add_subtitle(slide, "All recoverable. Knowing them up front saves the all-nighter.",
                 top=Inches(1.95))
    bugs = [
        ("All events arrive at once at the end",
         "A proxy or middleware is buffering. nginx: 'proxy_buffering off' or header 'X-Accel-Buffering: no'. Disable gzip on event-stream endpoints."),
        ("Connection dies after 60s of silence",
         "Idle-killer somewhere (LB, proxy, mobile network). Send ': ping\\n\\n' every 15-20 seconds."),
        ("Browser shows 'reconnecting' forever",
         "Server is returning a non-200 status. EventSource keeps retrying but never succeeds. Check server logs and CORS headers."),
        ("Authorization header doesn't work",
         "EventSource doesn't allow custom headers. Use cookies (best), short-lived URL token, or the event-source-polyfill library."),
        ("Sometimes I get duplicate events after reconnect",
         "Server replay is using >= instead of > the Last-Event-ID. Use strict greater-than."),
        ("Closing the tab leaves connection open server-side",
         "Normal - server only finds out at next heartbeat write. Clean up subscribers in your generator's 'finally' block."),
    ]
    top = Inches(2.55); h = Inches(0.72); gap = Inches(0.06)
    for i, (sym, fix) in enumerate(bugs):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        sx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.08), Inches(11.6), Inches(0.3))
        sp = sx.text_frame.paragraphs[0]
        sr = sp.add_run(); sr.text = sym
        sr.font.name = FONT_HEAD; sr.font.size = Pt(13); sr.font.bold = True; sr.font.color.rgb = WEBHOOK
        fx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.38), Inches(11.6), Inches(0.4))
        ff = fx.text_frame; ff.word_wrap = True
        fp = ff.paragraphs[0]
        fr = fp.add_run(); fr.text = fix
        fr.font.name = FONT_BODY; fr.font.size = Pt(11); fr.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ----- SSE section -----
def sse_section():
    section_divider("Server-Sent Events", SSE)

    # PAIR 1: SSE basics
    concept_slide("SSE", SSE,
                  "SSE - the wire format",
                  "One HTTP connection. Server streams labelled events; client reads as they arrive.",
                  DIAG / "08_sse_basic.png")
    example_slide("SSE", SSE,
                  "Raj watching his order status live",
                  "Server pushes 'paid -> cooking -> out for delivery -> delivered' as events. Zero polling.",
                  PERSONAS / "p_sse_order.png")

    # PAIR 2: reconnect / LLM streaming as a follow-on example
    concept_slide("SSE", SSE,
                  "Auto-reconnect with Last-Event-ID",
                  "Connection drops? Browser reconnects and tells the server what it last saw. Server resumes.",
                  DIAG / "09_sse_reconnect.png")
    example_slide("SSE", SSE,
                  "Streaming an LLM response (the ChatGPT pattern)",
                  "Backend proxies OpenAI's SSE stream as SSE to the browser. Same wire format both sides.",
                  PERSONAS / "p_sse_llm.png")

    # Deep dives
    slide_sse_grammar()
    slide_sse_eventsource_api()
    slide_sse_server_code()
    slide_sse_must_should()
    slide_sse_cost_picture()
    slide_sse_gotchas()

    # Where you've seen it
    slide = add_blank_slide()
    add_pattern_badge(slide, "SSE", SSE)
    add_title(slide, "You've used SSE thousands of times", top=Inches(1.15))
    add_subtitle(slide, "Anywhere you've watched text appear word-by-word in 2023+.", top=Inches(1.95))
    items = [
        "ChatGPT, Claude.ai, Gemini chat - the typing animation IS SSE",
        "Cursor, GitHub Copilot Chat, Continue.dev - same pattern",
        "OpenAI / Anthropic / Mistral SDKs with stream=True - SSE clients",
        "Vercel AI SDK useChat hook - SSE client + parser",
        "Vercel / Render / Netlify deploy logs streaming in browser - SSE",
        "MCP servers (Claude Desktop tools) - HTTP+SSE / Streamable HTTP transport",
        "LangSmith and LangFuse trace viewers - SSE",
    ]
    add_bullets(slide, items, top=Inches(2.7), size=16)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ws_scaling():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBSOCKET", WS)
    add_title(slide, "WebSocket scaling - what changes at each tier", top=Inches(1.15))
    add_subtitle(slide, "Each level adds infrastructure. Know which one you're at.", top=Inches(1.95))
    tiers = [
        ("Up to ~100", "Anything works", "One process. Default config. Pick by ergonomics."),
        ("100 - 10K", "Tune the basics", "Async I/O required. ulimit -n 65535. Heartbeat every 20s. LB idle timeout to 5 min."),
        ("10K - 100K", "Pub-sub backbone", "Multiple processes/instances. Redis or NATS pub-sub for cross-server broadcasts. Sticky LB sessions. Connection metrics."),
        ("100K - 1M", "Specialised tools", "Phoenix Channels (Elixir), Centrifugo, Ably, Pusher, PartyKit. Kernel tuning. Separate WS gateway tier from app servers. Edge delivery (Cloudflare/Fastly)."),
        ("Over 1M", "Specialised team", "You're past blog-post territory. Hire someone who's done it before. This is a full ops surface."),
    ]
    top = Inches(2.6); h = Inches(0.78); gap = Inches(0.08)
    for i, (tier, headline, body) in enumerate(tiers):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.0), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        tx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.12), Inches(2.3), Inches(0.4))
        tp = tx.text_frame.paragraphs[0]
        tr = tp.add_run(); tr.text = tier
        tr.font.name = FONT_HEAD; tr.font.size = Pt(14); tr.font.bold = True; tr.font.color.rgb = WS
        hx = slide.shapes.add_textbox(Inches(3.2), y + Inches(0.12), Inches(3.0), Inches(0.4))
        hp = hx.text_frame.paragraphs[0]
        hr = hp.add_run(); hr.text = headline
        hr.font.name = FONT_HEAD; hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = ACCENT
        bx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.4), Inches(11.5), Inches(0.4))
        bp = bx.text_frame.paragraphs[0]
        br = bp.add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(12); br.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ws_gotchas():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBSOCKET", WS)
    add_title(slide, "WebSocket bugs you'll hit, and how to ship around them", top=Inches(1.15))
    add_subtitle(slide, "Most production WS pain comes from these.",
                 top=Inches(1.95))
    bugs = [
        ("No reconnect logic on the client",
         "Network blips silently kill your client forever. Always write a reconnecting wrapper with exponential backoff + jitter."),
        ("No heartbeat",
         "Half-open TCP (network died but TCP doesn't know) looks alive until you write. Send {type:'ping'} every 20-30s; force reconnect if no pong."),
        ("Server is sync (Flask sync, Django pre-Channels)",
         "Exhausts worker pool at ~50 connections. WS needs async (asyncio/uvicorn, Node, Go) - period."),
        ("Custom auth via browser headers",
         "Browser WebSocket API does NOT allow custom headers. Use cookies, a subprotocol-as-token, or a short-lived URL token (ws://...?ticket=xxx)."),
        ("No backpressure handling",
         "Slow client + fast server = OOM. Bound your send queues; close clients that fall behind."),
        ("Works on dev (one server) breaks on prod (LB + multiple servers)",
         "Sticky sessions on LB; pub-sub backbone (Redis) so broadcasts cross processes; check LB actually supports Upgrade header."),
    ]
    top = Inches(2.55); h = Inches(0.72); gap = Inches(0.06)
    for i, (sym, fix) in enumerate(bugs):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        sx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.08), Inches(11.6), Inches(0.3))
        sp = sx.text_frame.paragraphs[0]
        sr = sp.add_run(); sr.text = sym
        sr.font.name = FONT_HEAD; sr.font.size = Pt(13); sr.font.bold = True; sr.font.color.rgb = WEBHOOK
        fx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.38), Inches(11.6), Inches(0.4))
        ff = fx.text_frame; ff.word_wrap = True
        fp = ff.paragraphs[0]
        fr = fp.add_run(); fr.text = fix
        fr.font.name = FONT_BODY; fr.font.size = Pt(11); fr.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ----- WebSocket section -----
def ws_section():
    section_divider("WebSockets", WS)

    # PAIR 1: handshake -> chat
    concept_slide("WEBSOCKET", WS,
                  "WebSocket - the handshake",
                  "Starts as HTTP. Upgrades to a custom frame protocol on the same TCP. Both sides talk freely.",
                  DIAG / "10_websocket_handshake.png")
    example_slide("WEBSOCKET", WS,
                  "Raj and Sam chat about the buzzer",
                  "Both can send anytime. Typing indicators are free. One TCP connection per chat.",
                  PERSONAS / "p_ws_chat.png")

    slide_ws_handshake_wire()
    slide_ws_frames()

    # PAIR 2: broadcast topology -> backend in real time
    concept_slide("WEBSOCKET", WS,
                  "Broadcast topology",
                  "Server tracks connected clients and fans messages out. Across processes, use Redis pub-sub.",
                  DIAG / "11_websocket_chat.png")

    slide_ws_chat_protocol()
    slide_ws_reconnecting()
    slide_ws_pubsub_cross_server()

    # Deep dives
    slide_ws_scaling()
    slide_ws_gotchas()

    # SSE vs WS comparison
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBSOCKET", WS)
    add_title(slide, "SSE vs WebSocket - when each one wins", top=Inches(1.15))
    add_subtitle(slide, "Default to SSE. WebSocket only when you NEED bidirectional.", top=Inches(1.95))
    # comparison table-like layout
    headers = ["", "SSE", "WebSocket"]
    rows = [
        ("Direction",        "Server -> client only",         "Both ways"),
        ("Protocol",         "Plain HTTP",                    "HTTP upgrade to ws://"),
        ("Browser API",      "EventSource (built in)",        "WebSocket (built in)"),
        ("Auto-reconnect",   "Yes, free",                     "No, you write it"),
        ("Binary support",   "No (text only)",                "Yes"),
        ("Through proxies",  "Almost always works",           "Sometimes blocked"),
        ("Best for",         "LLM tokens, status, logs",      "Chat, voice, games, collab"),
    ]
    top = Inches(2.7); row_h = Inches(0.45)
    col_x = [Inches(0.6), Inches(4.5), Inches(8.5)]
    col_w = [Inches(3.8), Inches(3.9), Inches(4.7)]
    # header
    for i, h in enumerate(headers):
        hx = slide.shapes.add_textbox(col_x[i], top, col_w[i], Inches(0.4))
        hf = hx.text_frame; hp = hf.paragraphs[0]
        hr = hp.add_run(); hr.text = h
        hr.font.name = FONT_HEAD; hr.font.size = Pt(15); hr.font.bold = True
        hr.font.color.rgb = ACCENT if i == 0 else (SSE if i == 1 else WS)
    # data rows
    for ri, row in enumerate(rows):
        y = top + Inches(0.5) + row_h * ri
        for ci, cell in enumerate(row):
            tx = slide.shapes.add_textbox(col_x[ci], y, col_w[ci], Inches(0.4))
            tf = tx.text_frame; tp = tf.paragraphs[0]
            tr = tp.add_run(); tr.text = cell
            tr.font.name = FONT_BODY; tr.font.size = Pt(13)
            tr.font.color.rgb = MUTED if ci == 0 else TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ----- Decision matrix section -----
def decision_section():
    section_divider("Decision Matrix", ACCENT)

    # PAIR: generic decision tree vs Maya's framing
    concept_slide("DECIDE", ACCENT,
                  "The decision tree - generic",
                  "Walk top to bottom. The first leaf that matches your data flow is the answer.",
                  DIAG / "12_decision_matrix.png")
    example_slide("DECIDE", ACCENT,
                  "Maya picking the pattern for a new feature",
                  "Same tree, framed as Maya's design checklist. Same answer, more memorable.",
                  PERSONAS / "p_decision_with_persona.png")

    slide_cost_ranking()
    slide_recipes_extended()
    slide_scale_tiers()


def slide_ai_long_task():
    slide = add_blank_slide()
    add_pattern_badge(slide, "AI SCENARIO", LLM)
    add_title(slide, "AI scenario - the agent kicks off a long task", top=Inches(1.15))
    add_subtitle(slide, "Research agent takes 3-5 min. Three valid patterns to surface progress.",
                 top=Inches(1.95))
    opts = [
        ("Polling",   POLLING,
         "POST /jobs returns {id}. UI polls /jobs/{id} every 5s. Spinner -> done toast. Simplest, no progress detail."),
        ("SSE for progress",  SSE,
         "POST /jobs returns {id}. UI opens /jobs/{id}/stream. Server pushes 'searching web', 'reading 3 docs', 'writing'. Best UX - feels fast even when slow."),
        ("Webhook callback",  WEBHOOK,
         "When the worker is external (OpenAI Batch, your own remote service). Pass callback_url; remote service POSTs when done. Best for fully async."),
    ]
    top = Inches(2.7); h = Inches(1.0); gap = Inches(0.1)
    for i, (name, color, body) in enumerate(opts):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.0), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = color; card.line.width = Pt(1.5); card.shadow.inherit = False
        nx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.18), Inches(3.0), Inches(0.4))
        np = nx.text_frame.paragraphs[0]
        nr = np.add_run(); nr.text = name
        nr.font.name = FONT_HEAD; nr.font.size = Pt(16); nr.font.bold = True; nr.font.color.rgb = color
        bx = slide.shapes.add_textbox(Inches(3.7), y + Inches(0.18), Inches(8.7), Inches(0.7))
        bf = bx.text_frame; bf.word_wrap = True
        bp = bf.paragraphs[0]
        br = bp.add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(13); br.font.color.rgb = TEXT
    add_text(slide, "Rule of thumb: SSE for progress within seconds. Webhook callback when the wait is hours.",
             top=Inches(6.4), size=14, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ai_voice_agent():
    slide = add_blank_slide()
    add_pattern_badge(slide, "AI SCENARIO", LLM)
    add_title(slide, "AI scenario - voice or interruptible streaming", top=Inches(1.15))
    add_subtitle(slide, "Microphone in, speaker out, mid-stream interrupt. Only WebSocket fits.",
                 top=Inches(1.95))
    # Visual: simple flow diagram in native PPT shapes
    # Boxes: mic, backend, ASR, LLM, TTS, speaker
    boxes = [
        ("Mic",      Inches(0.7),  Inches(3.1), WS),
        ("Backend",  Inches(2.6),  Inches(3.1), ACCENT),
        ("ASR",      Inches(4.5),  Inches(2.5), LLM),
        ("LLM",      Inches(4.5),  Inches(3.7), LLM),
        ("TTS",      Inches(6.4),  Inches(3.1), LLM),
        ("Speaker",  Inches(8.3),  Inches(3.1), WS),
    ]
    for label, x, y, color in boxes:
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.5), Inches(0.8))
        b.fill.solid(); b.fill.fore_color.rgb = BG_PANEL
        b.line.color.rgb = color; b.line.width = Pt(1.5); b.shadow.inherit = False
        tx = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(1.5), Inches(0.4))
        tp = tx.text_frame.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run(); tr.text = label
        tr.font.name = FONT_HEAD; tr.font.size = Pt(14); tr.font.bold = True; tr.font.color.rgb = color
    # Label across the middle
    cap = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.5))
    cf = cap.text_frame; cf.word_wrap = True
    cp = cf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = "Audio frames flow both ways over one WebSocket  ·  interrupting sends a 'stop' upstream while audio still streams down"
    cr.font.name = FONT_BODY; cr.font.size = Pt(13); cr.font.italic = True; cr.font.color.rgb = MUTED
    # Why others fail
    add_text(slide, "\n".join([
        "SSE   one-way (server -> client). Mic input has nowhere to go.",
        "Polling   latency is way too high; users would hear stutters.",
        "Webhooks   server-to-server only; no path from a user's microphone.",
        "WebSocket   binary + bidirectional + low overhead.  Used by OpenAI Realtime, ElevenLabs, Vapi, Retell.",
    ]), top=Inches(5.8), size=13, color=TEXT)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ----- Composing all 4 -----
def composition_section():
    section_divider("All Four, Together", LLM)

    # PAIR: generic AI app architecture vs LiveOrder
    concept_slide("COMPOSE", LLM,
                  "A real AI app uses all four",
                  "Generic architecture. Each edge labelled with which pattern carries it.",
                  DIAG / "13_ai_app_all_patterns.png")
    example_slide("COMPOSE", LLM,
                  "LiveOrder - all 4 in Raj's order journey",
                  "Same composition, our personas. This is Project 3 in the repo.",
                  PERSONAS / "p_liveorder_full.png")

    # PAIR: generic MCP vs Maya's MCP
    concept_slide("COMPOSE", LLM,
                  "MCP - the transport",
                  "Claude calls tools over HTTP+SSE / Streamable HTTP. One request -> many progress events.",
                  DIAG / "14_mcp_architecture.png")
    example_slide("COMPOSE", LLM,
                  "Maya exposing LiveOrder data to Claude",
                  "Maya's MCP server lets Claude query LiveOrder. The transport is SSE you've already learned.",
                  PERSONAS / "p_mcp_persona.png")


# ----- Workshop materials -----
def materials_section():
    section_divider("Workshop Kit", ACCENT)

    # The 3 projects
    slide = add_blank_slide()
    add_title(slide, "Three projects to fork")
    add_subtitle(slide, "Each runs locally. Each compares or composes patterns differently.")
    projects = [
        ("Project 1", "Streaming Chat - 3 ways",
         "Same LLM chat exposed via polling, SSE, and WebSocket side-by-side. Feel the UX difference."),
        ("Project 2", "Webhook Dashboard",
         "Stripe-style webhook intake (HMAC, dedup, fast 200) + receiving-side polling vs SSE comparison."),
        ("Project 3", "LiveOrder Capstone",
         "All 4 patterns in one mini-app: order tracking (SSE), driver chat (WS), payment (webhook), revenue report (polling), AI recommender (SSE/LLM)."),
    ]
    top = Inches(2.5); h = Inches(1.3); gap = Inches(0.2)
    for i, (num, name, body) in enumerate(projects):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.0), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        # left: project number
        nx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.25), Inches(1.6), Inches(0.5))
        np = nx.text_frame.paragraphs[0]
        nr = np.add_run(); nr.text = num
        nr.font.name = FONT_HEAD; nr.font.size = Pt(18); nr.font.bold = True; nr.font.color.rgb = ACCENT
        # middle: project name
        mx = slide.shapes.add_textbox(Inches(2.6), y + Inches(0.2), Inches(9.6), Inches(0.5))
        mp = mx.text_frame.paragraphs[0]
        mr = mp.add_run(); mr.text = name
        mr.font.name = FONT_HEAD; mr.font.size = Pt(20); mr.font.bold = True; mr.font.color.rgb = TEXT
        # bottom: description
        bx = slide.shapes.add_textbox(Inches(2.6), y + Inches(0.7), Inches(9.6), Inches(0.6))
        bf = bx.text_frame; bf.word_wrap = True
        bp = bf.paragraphs[0]
        br = bp.add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(13); br.font.color.rgb = MUTED
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))

    # Examples + Postman
    slide = add_blank_slide()
    add_title(slide, "Lab folders + Postman collection")
    add_subtitle(slide, "Read the code, run it, click Send, see the wire.")
    items = [
        "examples/01..07 - one tiny self-contained folder per topic with server.py + client.py",
        "examples/qa.sh - runs every example end-to-end (49 assertions, all green)",
        "postman/ - 53 requests in 9 folders, with pre-request scripts for HMAC + JWT",
        "diagrams/ - 14 technical + 10 persona scenario diagrams, all generated from Python builders",
        "concepts/ - 8 deep-dive markdown docs (~2900 lines) as the workbook companion",
        "TEACHING.md - the live-class playbook with talking points per slide",
    ]
    add_bullets(slide, items, top=Inches(2.5), size=15)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))

    # Repo
    slide = add_blank_slide()
    add_title(slide, "Get the kit")
    add_subtitle(slide, "Clone, run the QA, fork the project that fits what you're building.")
    add_text(slide, "github.com/fnusatvik07/real-time-events",
             top=Inches(2.5), size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "\n".join([
        "git clone https://github.com/fnusatvik07/real-time-events",
        "cd real-time-events",
        "python -m venv .venv && source .venv/bin/activate",
        "uv pip install fastapi 'uvicorn[standard]' httpx websockets python-dotenv openai",
        "cd examples && bash qa.sh    # 49 checks, ~2 minutes",
    ]), left=Inches(2.0), top=Inches(3.5), width=Inches(9.0), height=Inches(2.5),
        size=15, color=TEXT, mono=True)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_recap():
    slide = add_blank_slide()
    add_title(slide, "What to remember")
    add_subtitle(slide, "Print this on a sticky note.")
    # 4 cards, one per pattern, each with its 1-liner
    col_w = Inches(2.9); col_h = Inches(3.6); gap = Inches(0.2); start_x = Inches(0.7)
    top = Inches(2.7)
    cards = [
        ("POLLING",    POLLING, "Slow checks, rare events.\nOK to be a few seconds late."),
        ("WEBHOOKS",   WEBHOOK, "External trigger.\nSign + dedup + return 200 fast."),
        ("SSE",        SSE,     "Server -> client streaming.\nDefault for live UI."),
        ("WEBSOCKETS", WS,      "Both ways or binary only.\nReach last."),
    ]
    for i, (name, color, body) in enumerate(cards):
        x = start_x + (col_w + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = color; card.line.width = Pt(2); card.shadow.inherit = False
        tx = slide.shapes.add_textbox(x, top + Inches(0.5), col_w, Inches(0.6))
        tp = tx.text_frame.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run(); tr.text = name
        tr.font.name = FONT_HEAD; tr.font.size = Pt(20); tr.font.bold = True; tr.font.color.rgb = color
        bx = slide.shapes.add_textbox(x + Inches(0.2), top + Inches(1.4), col_w - Inches(0.4), Inches(2.0))
        bf = bx.text_frame; bf.word_wrap = True
        for j, line in enumerate(body.split("\n")):
            bp = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run(); br.text = line
            br.font.name = FONT_BODY; br.font.size = Pt(14); br.font.color.rgb = TEXT
    add_text(slide,
             "Default SSE. WebSocket only when bidirectional. Polling for slow. Webhooks for external.",
             top=Inches(6.5), size=15, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_qa():
    slide = add_blank_slide()
    # Big "Q&A" centered
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12.0), Inches(2.0))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Questions?"
    r.font.name = FONT_HEAD; r.font.size = Pt(96); r.font.bold = True; r.font.color.rgb = ACCENT
    # subtitle
    sx = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(12.0), Inches(0.6))
    sf = sx.text_frame
    sp = sf.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run(); sr.text = "Bring your own use case. We'll map it to a pattern."
    sr.font.name = FONT_BODY; sr.font.size = Pt(22); sr.font.italic = True; sr.font.color.rgb = MUTED
    # github at bottom
    gx = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.5))
    gf = gx.text_frame
    gp = gf.paragraphs[0]; gp.alignment = PP_ALIGN.CENTER
    gr = gp.add_run(); gr.text = "github.com/fnusatvik07/real-time-events"
    gr.font.name = FONT_CODE; gr.font.size = Pt(18); gr.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ===========================================================================
# Deep-dive slides covering the rest of the concept md material
# ===========================================================================

# ---- Overview / opening -----
def slide_liveorder_feature_map():
    slide = add_blank_slide()
    add_title(slide, "LiveOrder: one app, all four patterns")
    add_subtitle(slide, "Same backend Maya runs. Each feature picks the pattern that fits.")
    rows = [
        ("Stripe says Raj paid",            "External event",       "WEBHOOK",  WEBHOOK),
        ("Priya's tablet shows new orders", "Server -> restaurant", "SSE",      SSE),
        ("Raj's order status timeline",     "Server -> customer",   "SSE",      SSE),
        ("Raj chats with Sam about buzzer", "Two-way live",         "WEBSOCKET", WS),
        ("Raj sees Sam's scooter on map",   "Frequent updates",     "SSE / fast polling", SSE),
        ("Maya checks nightly pricing job", "Rare batch check",     "POLLING",   POLLING),
    ]
    headers = ("FEATURE", "DATA FLOW", "PATTERN")
    col_w = [Inches(5.0), Inches(4.0), Inches(3.3)]
    top = Inches(2.4); row_h = Inches(0.55)
    x = Inches(0.6)
    for i, h in enumerate(headers):
        hx = slide.shapes.add_textbox(x, top, col_w[i], Inches(0.35))
        hr = hx.text_frame.paragraphs[0].add_run(); hr.text = h
        hr.font.name = FONT_HEAD; hr.font.size = Pt(12); hr.font.bold = True; hr.font.color.rgb = ACCENT
        x += col_w[i]
    for ri, (feat, flow, pat, col) in enumerate(rows):
        y = top + Inches(0.45) + row_h * ri
        # subtle alternating panel
        if ri % 2 == 0:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(0.55), y - Inches(0.05),
                                          Inches(12.3), row_h)
            band.fill.solid(); band.fill.fore_color.rgb = BG_PANEL
            band.line.fill.background(); band.shadow.inherit = False
        x = Inches(0.6)
        for ci, cell in enumerate((feat, flow, pat)):
            tx = slide.shapes.add_textbox(x, y, col_w[ci], row_h - Inches(0.1))
            tp = tx.text_frame.paragraphs[0]
            tr = tp.add_run(); tr.text = cell
            tr.font.name = FONT_BODY; tr.font.size = Pt(14)
            tr.font.color.rgb = col if ci == 2 else TEXT
            tr.font.bold = (ci == 2)
            x += col_w[ci]
    add_text(slide, "Same backend. Four patterns. Each one chosen because it fits that feature's shape.",
             top=Inches(6.4), size=14, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ---- HTTP fundamentals -----
def slide_http_request_response_raw():
    slide = add_blank_slide()
    add_title(slide, "HTTP on the wire - it's just text")
    add_subtitle(slide, "Request and response are both plain text with a tiny structure. That's the entire protocol.")
    # Two code panels side by side
    panel_w = Inches(6.0); panel_h = Inches(4.2); top = Inches(2.4)
    # Left - REQUEST
    lx = Inches(0.6)
    lcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, top, panel_w, panel_h)
    lcard.fill.solid(); lcard.fill.fore_color.rgb = BG_PANEL
    lcard.line.color.rgb = ACCENT; lcard.line.width = Pt(1.0); lcard.shadow.inherit = False
    lh = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.15), panel_w - Inches(0.5), Inches(0.4))
    lhr = lh.text_frame.paragraphs[0].add_run(); lhr.text = "REQUEST  (what your browser sends)"
    lhr.font.name = FONT_HEAD; lhr.font.size = Pt(13); lhr.font.bold = True; lhr.font.color.rgb = ACCENT
    req = """GET /api/users HTTP/1.1
Host: api.example.com
Accept: application/json
Authorization: Bearer eyJhbGciOiJ...
User-Agent: Mozilla/5.0 ...

"""
    tx = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.65), panel_w - Inches(0.5), panel_h - Inches(0.9))
    tf = tx.text_frame; tf.word_wrap = True
    for i, line in enumerate(req.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_CODE; r.font.size = Pt(12); r.font.color.rgb = TEXT
    # Right - RESPONSE
    rx = Inches(6.8)
    rcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, top, panel_w, panel_h)
    rcard.fill.solid(); rcard.fill.fore_color.rgb = BG_PANEL
    rcard.line.color.rgb = SSE; rcard.line.width = Pt(1.0); rcard.shadow.inherit = False
    rh = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.15), panel_w - Inches(0.5), Inches(0.4))
    rhr = rh.text_frame.paragraphs[0].add_run(); rhr.text = "RESPONSE  (what the server replies)"
    rhr.font.name = FONT_HEAD; rhr.font.size = Pt(13); rhr.font.bold = True; rhr.font.color.rgb = SSE
    res = """HTTP/1.1 200 OK
Content-Type: application/json
Content-Length: 348
Cache-Control: no-store

{"users":[
  {"id":1,"name":"Maya"},
  {"id":2,"name":"Raj"}
]}"""
    tx = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.65), panel_w - Inches(0.5), panel_h - Inches(0.9))
    tf = tx.text_frame; tf.word_wrap = True
    for i, line in enumerate(res.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_CODE; r.font.size = Pt(12); r.font.color.rgb = TEXT
    add_text(slide, "Method + path + headers + blank line + optional body.  Status line + headers + blank line + body.",
             top=Inches(6.85), size=12, color=MUTED, align=PP_ALIGN.CENTER, italic_via=True)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_stateless_concept():
    slide = add_blank_slide()
    add_title(slide, "Stateless - the single property that makes real-time hard")
    add_subtitle(slide, "The server has no memory between requests. The client carries everything.")
    body = (
        "Two requests one second apart  =  two strangers walking up to a counter.\n"
        "The server doesn't remember 'Maya is logged in' - it re-derives it from the token Maya's browser\n"
        "sends with every single request.\n\n"
        "This is GREAT for scale (any server can serve any request) but TERRIBLE for real-time:\n\n"
        "When a new message arrives on the server for Raj, the server cannot say 'let me notify Raj' -\n"
        "there is no open line to Raj. There is no phone number. There is no email address baked in.\n\n"
        "The four real-time patterns are all answers to the question:\n"
        "  'how do we work around HTTP being request/response, stateless, and client-initiated?'"
    )
    add_text(slide, body, top=Inches(2.4), size=15, color=TEXT)
    # Three workarounds at the bottom
    items = [
        ("Polling",   "client keeps asking",                          POLLING),
        ("SSE / WS",  "client opens a long-lived line; server pushes", SSE),
        ("Webhook",   "an external system already had a line to your server", WEBHOOK),
    ]
    top = Inches(5.8); w = Inches(4.0); gap = Inches(0.15)
    for i, (head, body, color) in enumerate(items):
        x = Inches(0.6) + (w + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, Inches(1.0))
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = color; card.line.width = Pt(1.5); card.shadow.inherit = False
        hx = slide.shapes.add_textbox(x + Inches(0.2), top + Inches(0.1), w - Inches(0.4), Inches(0.35))
        hr = hx.text_frame.paragraphs[0].add_run(); hr.text = head
        hr.font.name = FONT_HEAD; hr.font.size = Pt(14); hr.font.bold = True; hr.font.color.rgb = color
        bx = slide.shapes.add_textbox(x + Inches(0.2), top + Inches(0.45), w - Inches(0.4), Inches(0.5))
        bf = bx.text_frame; bf.word_wrap = True
        br = bf.paragraphs[0].add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(12); br.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_http_versions():
    slide = add_blank_slide()
    add_title(slide, "HTTP versions in 90 seconds")
    add_subtitle(slide, "All four real-time patterns are identical across versions. Only scaling limits differ.")
    rows = [
        ("HTTP/1.0", "one TCP connection per request",
         "mostly historical"),
        ("HTTP/1.1", "keep-alive (re-use connections), chunked transfer",
         "what most 'long-lived' patterns assume; SSE works here but is limited to 6 connections per host in browsers"),
        ("HTTP/2",   "one TCP connection multiplexes many parallel requests",
         "SSE works much better here; modern default everywhere"),
        ("HTTP/3",   "same idea as HTTP/2 but runs on UDP (QUIC)",
         "better on flaky mobile; connection survives IP changes"),
    ]
    top = Inches(2.5); h = Inches(0.95); gap = Inches(0.1)
    for i, (ver, change, care) in enumerate(rows):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        vx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.15), Inches(2.0), Inches(0.4))
        vr = vx.text_frame.paragraphs[0].add_run(); vr.text = ver
        vr.font.name = FONT_HEAD; vr.font.size = Pt(15); vr.font.bold = True; vr.font.color.rgb = ACCENT
        cx = slide.shapes.add_textbox(Inches(3.0), y + Inches(0.15), Inches(4.5), Inches(0.4))
        cr = cx.text_frame.paragraphs[0].add_run(); cr.text = change
        cr.font.name = FONT_BODY; cr.font.size = Pt(13); cr.font.color.rgb = TEXT; cr.font.bold = True
        wx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.5), Inches(11.5), Inches(0.4))
        wf = wx.text_frame; wf.word_wrap = True
        wr = wf.paragraphs[0].add_run(); wr.text = "why care: " + care
        wr.font.name = FONT_BODY; wr.font.size = Pt(12); wr.font.color.rgb = MUTED
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ---- Polling additions -----
def slide_polling_code():
    slide = add_blank_slide()
    add_pattern_badge(slide, "POLLING", POLLING)
    add_title(slide, "Short polling in 18 lines", top=Inches(1.15))
    add_subtitle(slide, "The whole pattern - client setInterval, server returns current state.",
                 top=Inches(1.95))
    panel_w = Inches(6.0); panel_h = Inches(4.4); top = Inches(2.5)
    # Client JS
    lx = Inches(0.6)
    lcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, top, panel_w, panel_h)
    lcard.fill.solid(); lcard.fill.fore_color.rgb = BG_PANEL
    lcard.line.color.rgb = POLLING; lcard.line.width = Pt(1.0); lcard.shadow.inherit = False
    lh = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.15), panel_w - Inches(0.5), Inches(0.4))
    lhr = lh.text_frame.paragraphs[0].add_run(); lhr.text = "client.js  (8 lines)"
    lhr.font.name = FONT_HEAD; lhr.font.size = Pt(13); lhr.font.bold = True; lhr.font.color.rgb = POLLING
    code_js = """let last = null;
const id = setInterval(async () => {
  const r = await fetch(`/orders/${o}/status`);
  const { status } = await r.json();
  if (status !== last) {
    updateUI(status);
    last = status;
  }
  if (status === "delivered") clearInterval(id);
}, 3000);   // poll every 3 seconds
"""
    tx = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.65), panel_w - Inches(0.5), panel_h - Inches(0.9))
    tf = tx.text_frame; tf.word_wrap = False
    for i, line in enumerate(code_js.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_CODE; r.font.size = Pt(13); r.font.color.rgb = TEXT
    # Server Python
    rx = Inches(6.8)
    rcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, top, panel_w, panel_h)
    rcard.fill.solid(); rcard.fill.fore_color.rgb = BG_PANEL
    rcard.line.color.rgb = SSE; rcard.line.width = Pt(1.0); rcard.shadow.inherit = False
    rh = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.15), panel_w - Inches(0.5), Inches(0.4))
    rhr = rh.text_frame.paragraphs[0].add_run(); rhr.text = "server.py  (4 lines)"
    rhr.font.name = FONT_HEAD; rhr.font.size = Pt(13); rhr.font.bold = True; rhr.font.color.rgb = SSE
    code_py = """@app.get("/orders/{order_id}/status")
def get_status(order_id: int):
    order = db.get_order(order_id)
    return {"status": order.status}
"""
    tx = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.65), panel_w - Inches(0.5), panel_h - Inches(0.9))
    tf = tx.text_frame; tf.word_wrap = False
    for i, line in enumerate(code_py.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_CODE; r.font.size = Pt(13); r.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_polling_cost_math():
    slide = add_blank_slide()
    add_pattern_badge(slide, "POLLING", POLLING)
    add_title(slide, "The hidden cost of short polling", top=Inches(1.15))
    add_subtitle(slide, "Raj's order: 25 minutes, status changes 5 times, polled every 3 seconds.",
                 top=Inches(1.95))
    # Three big stat tiles
    stats = [
        ("500",   "polls fired",        TEXT),
        ("5",     "actually meaningful", SSE),
        ("99%",   "wasted",             WEBHOOK),
    ]
    top = Inches(2.7); w = Inches(4.0); gap = Inches(0.1)
    for i, (num, label, col) in enumerate(stats):
        x = Inches(0.6) + (w + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, Inches(1.6))
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = col; card.line.width = Pt(1.5); card.shadow.inherit = False
        nx = slide.shapes.add_textbox(x, top + Inches(0.2), w, Inches(0.8))
        np = nx.text_frame.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
        nr = np.add_run(); nr.text = num
        nr.font.name = FONT_HEAD; nr.font.size = Pt(56); nr.font.bold = True; nr.font.color.rgb = col
        lx = slide.shapes.add_textbox(x, top + Inches(1.05), w, Inches(0.5))
        lp = lx.text_frame.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run(); lr.text = label
        lr.font.name = FONT_BODY; lr.font.size = Pt(14); lr.font.color.rgb = MUTED
    # Multiply by scale
    add_text(slide,
             "Now multiply by 10,000 concurrent customers:\n"
             "5,000,000 HTTP requests per hour just to track order statuses.\n"
             "99% of them returned nothing useful.",
             top=Inches(4.7), size=18, color=TEXT, align=PP_ALIGN.CENTER)
    # Punch line
    bx = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.6))
    bp = bx.text_frame.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "Lower interval -> smoother UX, higher cost.    Raise interval -> cheaper, laggy.    You cannot win both."
    br.font.name = FONT_HEAD; br.font.size = Pt(15); br.font.bold = True; br.font.italic = True; br.font.color.rgb = ACCENT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_polling_when_right():
    slide = add_blank_slide()
    add_pattern_badge(slide, "POLLING", POLLING)
    add_title(slide, "When polling IS the right call", top=Inches(1.15))
    add_subtitle(slide, "It is fashionable to roll eyes at polling. Resist. Big companies use it for the right things.",
                 top=Inches(1.95))
    rows = [
        ("Talking to an external API you don't control",
         "OpenAI Batch / Fine-tuning, S3 Glacier, AWS Lambda async invokes - all polled because they don't push to you."),
        ("The thing you're checking changes rarely",
         "A nightly batch job. Polling every minute is 60 reqs; SSE would be overkill."),
        ("Cost of being a few seconds stale is zero",
         "A dashboard of yesterday's sales numbers. Nobody's watching it tick."),
        ("You have very few clients",
         "Internal admin tool used by 5 people. 500 polls/hour is rounding-noise."),
        ("Platform / network kills long-lived connections",
         "Corporate proxies, iOS background fetch, aggressive serverless idle timeouts."),
        ("You're shipping in three days, alone",
         "Polling is the only pattern you can build, test, deploy, and debug in an afternoon."),
    ]
    top = Inches(2.6); h = Inches(0.62); gap = Inches(0.08)
    for i, (sym, fix) in enumerate(rows):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = POLLING; card.line.width = Pt(1.0); card.shadow.inherit = False
        sx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.06), Inches(11.6), Inches(0.3))
        sp = sx.text_frame.paragraphs[0]
        sr = sp.add_run(); sr.text = sym
        sr.font.name = FONT_HEAD; sr.font.size = Pt(13); sr.font.bold = True; sr.font.color.rgb = POLLING
        fx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.34), Inches(11.6), Inches(0.3))
        ff = fx.text_frame; ff.word_wrap = True
        fp = ff.paragraphs[0]
        fr = fp.add_run(); fr.text = fix
        fr.font.name = FONT_BODY; fr.font.size = Pt(11); fr.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ---- Webhook additions -----
def slide_webhook_role_flip():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "The role flip - webhooks are just inverted polling", top=Inches(1.15))
    add_subtitle(slide, "Same HTTP. Same servers. Just swap who is the client and who is the server.",
                 top=Inches(1.95))
    # Two big panels comparing
    panel_w = Inches(6.0); panel_h = Inches(3.6); top = Inches(2.6)
    # POLLING
    lx = Inches(0.6)
    lcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, top, panel_w, panel_h)
    lcard.fill.solid(); lcard.fill.fore_color.rgb = BG_PANEL
    lcard.line.color.rgb = POLLING; lcard.line.width = Pt(1.5); lcard.shadow.inherit = False
    lh = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.2), panel_w - Inches(0.5), Inches(0.5))
    lhr = lh.text_frame.paragraphs[0].add_run(); lhr.text = "POLLING"
    lhr.font.name = FONT_HEAD; lhr.font.size = Pt(20); lhr.font.bold = True; lhr.font.color.rgb = POLLING
    add_text(slide,
             "YOU are the client.\n"
             "THEY are the server.\n\n"
             "You ASK them.\n"
             "  GET /payments/recent\n\n"
             "  GET /payments/recent\n"
             "  GET /payments/recent\n"
             "  GET /payments/recent\n"
             "  GET ...",
             left=lx + Inches(0.4), top=top + Inches(0.85),
             width=panel_w - Inches(0.6), height=panel_h - Inches(1.0),
             size=14, color=TEXT)
    # WEBHOOK
    rx = Inches(6.8)
    rcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, top, panel_w, panel_h)
    rcard.fill.solid(); rcard.fill.fore_color.rgb = BG_PANEL
    rcard.line.color.rgb = WEBHOOK; rcard.line.width = Pt(1.5); rcard.shadow.inherit = False
    rh = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.2), panel_w - Inches(0.5), Inches(0.5))
    rhr = rh.text_frame.paragraphs[0].add_run(); rhr.text = "WEBHOOK"
    rhr.font.name = FONT_HEAD; rhr.font.size = Pt(20); rhr.font.bold = True; rhr.font.color.rgb = WEBHOOK
    add_text(slide,
             "THEY are the client.\n"
             "YOU are the server.\n\n"
             "They POST you (when it happens).\n"
             "  POST /webhooks/stripe\n"
             "  { type: payment.succeeded, ... }\n\n"
             "  (silence ... silence ...)\n"
             "  POST /webhooks/stripe",
             left=rx + Inches(0.4), top=top + Inches(0.85),
             width=panel_w - Inches(0.6), height=panel_h - Inches(1.0),
             size=14, color=TEXT)
    add_text(slide, "That role flip is the entire idea behind webhooks.",
             top=Inches(6.45), size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_webhook_payload():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "Webhook payload anatomy", top=Inches(1.15))
    add_subtitle(slide, "Stripe shape, but every provider follows the same outline.",
                 top=Inches(1.95))
    code_block(slide, """{
  "id":          "evt_1NXY2bAB4tH...",        # unique event id - USE THIS FOR DEDUP
  "type":        "payment_intent.succeeded",  # what happened - your switch statement
  "created":     1700000000,                  # timestamp - useful for ordering
  "data": {
    "object": {                               # the resource at event-fire moment
      "id":        "pi_3NXY...",
      "amount":    5000,
      "currency":  "usd",
      "customer":  "cus_abc123",
      "metadata":  { "order_id": "order_789" }
    }
  },
  "livemode":    true,
  "api_version": "2023-10-16"                 # sender's API version
}""", top=Inches(2.5), height=Inches(4.0), size=12)
    # Bottom callout
    add_text(slide,
             "Don't assume metadata fields you didn't set are present. Defensive parsing pays off.",
             top=Inches(6.7), size=13, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_webhook_receiver_code():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "A production-shape Stripe receiver (FastAPI, 22 lines)", top=Inches(1.15))
    add_subtitle(slide, "Verify, dedup, enqueue, return 200. Everything else happens in a worker.",
                 top=Inches(1.95))
    code_block(slide, """import hmac, hashlib, json, os
from fastapi import FastAPI, Request, HTTPException

app = FastAPI()
SECRET = os.environ["STRIPE_WEBHOOK_SECRET"]
seen = set()                                  # in prod: Redis SET with TTL

def verify(body: bytes, header: str) -> bool:
    parts  = dict(p.split("=", 1) for p in header.split(","))
    signed = f"{parts['t']}.{body.decode()}".encode()
    expect = hmac.new(SECRET.encode(), signed, hashlib.sha256).hexdigest()
    return hmac.compare_digest(expect, parts["v1"])     # constant-time!

@app.post("/webhooks/stripe")
async def stripe_webhook(req: Request):
    body = await req.body()                   # RAW bytes, not parsed JSON
    if not verify(body, req.headers.get("stripe-signature", "")):
        raise HTTPException(401, "bad signature")
    event = json.loads(body)
    if event["id"] in seen:                   # dedup BEFORE enqueue
        return {"received": True, "duplicate": True}
    seen.add(event["id"])
    await queue.put(event)                    # the real work happens elsewhere
    return {"received": True}""", top=Inches(2.5), height=Inches(4.4), size=11)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_webhook_real_world():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBHOOK", WEBHOOK)
    add_title(slide, "Webhooks in the wild", top=Inches(1.15))
    add_subtitle(slide, "Once you know to look, you see them everywhere.",
                 top=Inches(1.95))
    items = [
        ("Stripe",      "payments, refunds, disputes, subscriptions"),
        ("GitHub",      "pull requests, pushes, issues, deployments - powers every CI system"),
        ("Slack Events API", "messages, channel events, reactions - how every chat bot works"),
        ("Twilio",      "incoming SMS / call events / delivery status"),
        ("Calendly / Cal.com", "meeting booked, cancelled, rescheduled"),
        ("Shopify",     "orders, inventory updates, abandoned carts"),
        ("Auth0 / Clerk / Supabase Auth", "user signed up, password changed, MFA enabled"),
        ("OpenAI Batch API", "when a batch job finishes (newer services do this)"),
        ("Your own service",  "any time customers ask 'can you POST when X happens?'"),
    ]
    top = Inches(2.55); h = Inches(0.46); gap = Inches(0.04)
    for i, (name, what) in enumerate(items):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        nx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.08), Inches(3.5), Inches(0.3))
        nr = nx.text_frame.paragraphs[0].add_run(); nr.text = name
        nr.font.name = FONT_HEAD; nr.font.size = Pt(13); nr.font.bold = True; nr.font.color.rgb = WEBHOOK
        wx = slide.shapes.add_textbox(Inches(4.5), y + Inches(0.08), Inches(8.0), Inches(0.3))
        wr = wx.text_frame.paragraphs[0].add_run(); wr.text = what
        wr.font.name = FONT_BODY; wr.font.size = Pt(12); wr.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ---- SSE additions -----
def slide_sse_eventsource_api():
    slide = add_blank_slide()
    add_pattern_badge(slide, "SSE", SSE)
    add_title(slide, "EventSource - the browser's SSE client", top=Inches(1.15))
    add_subtitle(slide, "Three handlers, one close. The browser does parsing, reconnecting, Last-Event-ID.",
                 top=Inches(1.95))
    code_block(slide, """const es = new EventSource('/api/stream');

es.onmessage = (e) => console.log('default event:', e.data);

es.addEventListener('token', (e) => {                   // NAMED event
  const { text } = JSON.parse(e.data);
  appendToken(text);
});

es.onopen  = () => console.log('connected');
es.onerror = () => console.log('error or reconnecting');

es.close();   // otherwise it auto-reconnects forever
""", top=Inches(2.5), height=Inches(2.6), size=13)
    # Two cards: gives free / doesn't
    top = Inches(5.25); w = Inches(6.0); gap = Inches(0.13)
    for i, (head, color, items) in enumerate([
        ("FREE from EventSource", SSE,
         ["Auto-reconnect on dropped connection", "Last-Event-ID sent on reconnect", "Built into every modern browser"]),
        ("NOT included", WEBHOOK,
         ["Custom headers (no Authorization: Bearer ...)", "Binary frames (text only)", "Way to send data back (use fetch POST)"]),
    ]):
        x = Inches(0.6) + (w + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, Inches(1.7))
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = color; card.line.width = Pt(1.2); card.shadow.inherit = False
        hx = slide.shapes.add_textbox(x + Inches(0.2), top + Inches(0.1), w - Inches(0.4), Inches(0.35))
        hr = hx.text_frame.paragraphs[0].add_run(); hr.text = head
        hr.font.name = FONT_HEAD; hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = color
        for j, it in enumerate(items):
            ix = slide.shapes.add_textbox(x + Inches(0.2), top + Inches(0.45) + Inches(0.35) * j,
                                          w - Inches(0.4), Inches(0.3))
            ip = ix.text_frame.paragraphs[0]
            ir = ip.add_run(); ir.text = "·  " + it
            ir.font.name = FONT_BODY; ir.font.size = Pt(12); ir.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_sse_server_code():
    slide = add_blank_slide()
    add_pattern_badge(slide, "SSE", SSE)
    add_title(slide, "Minimal SSE endpoint (FastAPI)", top=Inches(1.15))
    add_subtitle(slide, "A generator yields events. StreamingResponse pipes them out.",
                 top=Inches(1.95))
    code_block(slide, """from fastapi import FastAPI
from fastapi.responses import StreamingResponse
import asyncio

app = FastAPI()

@app.get("/stream")
async def stream():
    async def gen():
        for i in range(10):
            yield f"data: chunk {i}\\n\\n"      # double-\\n ends each event
            await asyncio.sleep(1)
        yield "event: done\\ndata: complete\\n\\n"

    return StreamingResponse(gen(), media_type="text/event-stream", headers={
        "Cache-Control":     "no-cache, no-transform",
        "X-Accel-Buffering": "no",              # disable nginx buffering
    })""", top=Inches(2.5), height=Inches(4.4), size=13)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_sse_must_should():
    slide = add_blank_slide()
    add_pattern_badge(slide, "SSE", SSE)
    add_title(slide, "Three MUSTs and three SHOULDs for an SSE server", top=Inches(1.15))
    add_subtitle(slide, "Get these right and your SSE endpoint behaves the same in dev and prod.",
                 top=Inches(1.95))
    musts = [
        ("Set Content-Type: text/event-stream",
         "How the browser knows to keep the connection open and parse the body as events."),
        ("End every event with \\n\\n",
         "Blank line tells the client 'this event is complete'. Forget it and the client waits forever."),
        ("Disable proxy buffering",
         "X-Accel-Buffering: no for nginx; check LB docs for others. Without this, your stream batches at the proxy."),
    ]
    shoulds = [
        ("Send keep-alives",
         "Every 15-30s send ': ping\\n\\n'. Stops proxies/LB from dropping the idle line."),
        ("Cache-Control: no-cache, no-transform",
         "Stops any layer from caching the response."),
        ("Don't gzip the stream",
         "Compression buffers bytes to compress efficiently. Defeats streaming."),
    ]
    # Two columns
    panel_w = Inches(6.0); panel_h = Inches(4.5); top = Inches(2.55)
    for col_i, (head, color, items) in enumerate([
        ("MUST", SSE, musts),
        ("SHOULD", ACCENT, shoulds),
    ]):
        x = Inches(0.6) + (panel_w + Inches(0.13)) * col_i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, panel_w, panel_h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = color; card.line.width = Pt(1.5); card.shadow.inherit = False
        hx = slide.shapes.add_textbox(x + Inches(0.25), top + Inches(0.2), panel_w - Inches(0.5), Inches(0.45))
        hr = hx.text_frame.paragraphs[0].add_run(); hr.text = head
        hr.font.name = FONT_HEAD; hr.font.size = Pt(20); hr.font.bold = True; hr.font.color.rgb = color
        for j, (h2, body) in enumerate(items):
            y = top + Inches(0.8) + Inches(1.2) * j
            hx2 = slide.shapes.add_textbox(x + Inches(0.3), y, panel_w - Inches(0.5), Inches(0.35))
            hr2 = hx2.text_frame.paragraphs[0].add_run(); hr2.text = h2
            hr2.font.name = FONT_HEAD; hr2.font.size = Pt(13); hr2.font.bold = True; hr2.font.color.rgb = TEXT
            bx = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.35), panel_w - Inches(0.5), Inches(0.7))
            bf = bx.text_frame; bf.word_wrap = True
            br = bf.paragraphs[0].add_run(); br.text = body
            br.font.name = FONT_BODY; br.font.size = Pt(11); br.font.color.rgb = MUTED
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_sse_cost_picture():
    slide = add_blank_slide()
    add_pattern_badge(slide, "SSE", SSE)
    add_title(slide, "Maya picks SSE: who pays what", top=Inches(1.15))
    add_subtitle(slide, "Same 'track my order' feature. Numbers at 10K concurrent customers.",
                 top=Inches(1.95))
    # 4-row comparison
    rows = [
        ("Latency",         "0-3 seconds (poll gap)",      "near zero (push)"),
        ("Idle requests",   "10K every 3s = 3,300/sec",     "0 - heartbeat ping every 20s"),
        ("Server CPU",      "high - lots of work to say nothing", "low"),
        ("Memory",          "tiny - stateless",             "~1 GB at 10K (50-100 KB / open conn)"),
        ("User experience", "feels laggy on status changes", "instant"),
    ]
    headers = ("METRIC", "POLLING (every 3s)", "SSE")
    col_w = [Inches(3.2), Inches(4.8), Inches(4.1)]
    top = Inches(2.6); row_h = Inches(0.65)
    x = Inches(0.6)
    for i, h in enumerate(headers):
        hx = slide.shapes.add_textbox(x, top, col_w[i], Inches(0.4))
        hr = hx.text_frame.paragraphs[0].add_run(); hr.text = h
        hr.font.name = FONT_HEAD; hr.font.size = Pt(13); hr.font.bold = True
        hr.font.color.rgb = ACCENT if i == 0 else (POLLING if i == 1 else SSE)
        x += col_w[i]
    for ri, row in enumerate(rows):
        y = top + Inches(0.5) + row_h * ri
        x = Inches(0.6)
        for ci, cell in enumerate(row):
            tx = slide.shapes.add_textbox(x, y, col_w[ci], row_h - Inches(0.1))
            tf = tx.text_frame; tf.word_wrap = True
            tp = tf.paragraphs[0]
            tr = tp.add_run(); tr.text = cell
            tr.font.name = FONT_BODY; tr.font.size = Pt(13); tr.font.color.rgb = TEXT
            if ci == 0:
                tr.font.bold = True; tr.font.color.rgb = MUTED
            x += col_w[ci]
    add_text(slide,
             "10K open connections is the SSE sweet spot. Past 100K, plan for a pub-sub backbone.",
             top=Inches(6.6), size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ---- WebSocket additions -----
def slide_ws_handshake_wire():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBSOCKET", WS)
    add_title(slide, "The handshake on the wire", top=Inches(1.15))
    add_subtitle(slide, "Starts as HTTP. After 101 Switching Protocols, HTTP is gone - frames take over.",
                 top=Inches(1.95))
    panel_w = Inches(6.0); panel_h = Inches(3.6); top = Inches(2.55)
    # Request
    lx = Inches(0.6)
    lcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, top, panel_w, panel_h)
    lcard.fill.solid(); lcard.fill.fore_color.rgb = BG_PANEL
    lcard.line.color.rgb = WS; lcard.line.width = Pt(1.0); lcard.shadow.inherit = False
    lh = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.15), panel_w - Inches(0.5), Inches(0.4))
    lhr = lh.text_frame.paragraphs[0].add_run(); lhr.text = "Client REQUEST"
    lhr.font.name = FONT_HEAD; lhr.font.size = Pt(13); lhr.font.bold = True; lhr.font.color.rgb = WS
    req = """GET /chat HTTP/1.1
Host: api.liveorder.app
Upgrade: websocket
Connection: Upgrade
Sec-WebSocket-Key: dGhlIHNhbXBsZSBub25jZQ==
Sec-WebSocket-Version: 13"""
    tx = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.65), panel_w - Inches(0.5), panel_h - Inches(0.9))
    tf = tx.text_frame
    for i, line in enumerate(req.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_CODE; r.font.size = Pt(12); r.font.color.rgb = TEXT
    # Response
    rx = Inches(6.8)
    rcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, top, panel_w, panel_h)
    rcard.fill.solid(); rcard.fill.fore_color.rgb = BG_PANEL
    rcard.line.color.rgb = SSE; rcard.line.width = Pt(1.0); rcard.shadow.inherit = False
    rh = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.15), panel_w - Inches(0.5), Inches(0.4))
    rhr = rh.text_frame.paragraphs[0].add_run(); rhr.text = "Server RESPONSE"
    rhr.font.name = FONT_HEAD; rhr.font.size = Pt(13); rhr.font.bold = True; rhr.font.color.rgb = SSE
    res = """HTTP/1.1 101 Switching Protocols
Upgrade: websocket
Connection: Upgrade
Sec-WebSocket-Accept: s3pPLMBiTxaQ9kYGzzhZRbK+xOo=

[after this point: WebSocket frames,
 no more HTTP messages on this TCP socket]"""
    tx = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.65), panel_w - Inches(0.5), panel_h - Inches(0.9))
    tf = tx.text_frame
    for i, line in enumerate(res.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_CODE; r.font.size = Pt(12); r.font.color.rgb = TEXT
    add_text(slide,
             "Why the dance? Ports 80/443 are open everywhere; brand-new protocols on new ports get blocked.\nDress as HTTP, sneak through middleboxes, then upgrade.",
             top=Inches(6.4), size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ws_frames():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBSOCKET", WS)
    add_title(slide, "WebSocket frames - the unit of communication", top=Inches(1.15))
    add_subtitle(slide, "Five frame types. Libraries handle framing for you, but knowing helps debugging.",
                 top=Inches(1.95))
    rows = [
        ("Text",   "UTF-8 string. Most chat apps use these with JSON inside.",   WS),
        ("Binary", "arbitrary bytes. Audio, images, protobuf, msgpack.",         LLM),
        ("Ping",   "heartbeat. Library usually auto-replies with a Pong.",       ACCENT),
        ("Pong",   "reply to a Ping. Used to detect dead-but-not-yet-closed connections.", ACCENT),
        ("Close",  "graceful shutdown with status code (1000 normal, 1006 abnormal, 1011 server error).", WEBHOOK),
    ]
    top = Inches(2.55); h = Inches(0.7); gap = Inches(0.1)
    for i, (name, body, color) in enumerate(rows):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = color; card.line.width = Pt(1.0); card.shadow.inherit = False
        nx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.2), Inches(1.8), Inches(0.4))
        nr = nx.text_frame.paragraphs[0].add_run(); nr.text = name
        nr.font.name = FONT_HEAD; nr.font.size = Pt(15); nr.font.bold = True; nr.font.color.rgb = color
        bx = slide.shapes.add_textbox(Inches(2.8), y + Inches(0.2), Inches(9.6), Inches(0.4))
        bf = bx.text_frame; bf.word_wrap = True
        br = bf.paragraphs[0].add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(13); br.font.color.rgb = TEXT
    add_text(slide,
             "Each frame has: opcode + payload length + (client) mask key + payload bytes. Up to 2^63 bytes per message.",
             top=Inches(6.7), size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ws_reconnecting():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBSOCKET", WS)
    add_title(slide, "You WILL write ReconnectingWS", top=Inches(1.15))
    add_subtitle(slide, "Unlike EventSource, the WebSocket API gives you nothing for free. Here's the minimum.",
                 top=Inches(1.95))
    code_block(slide, """class ReconnectingWS {
  constructor(url) {
    this.url = url;
    this.backoff = 1000;            // start at 1s
    this.maxBackoff = 30000;
    this.connect();
  }
  connect() {
    this.ws = new WebSocket(this.url);
    this.ws.onopen    = () => { this.backoff = 1000; this.onopen?.(); };       // reset on success
    this.ws.onmessage = (e) => this.onmessage?.(e);
    this.ws.onclose   = () => {
      setTimeout(() => this.connect(), this.backoff + Math.random() * 500);    // + JITTER
      this.backoff = Math.min(this.backoff * 2, this.maxBackoff);              // exponential
    };
    this.ws.onerror   = () => this.ws.close();   // trigger onclose path
  }
  send(msg) {
    if (this.ws.readyState === WebSocket.OPEN) this.ws.send(msg);
    else this.queue.push(msg);      // optional: buffer until reconnect
  }
}""", top=Inches(2.5), height=Inches(4.4), size=11)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ws_chat_protocol():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBSOCKET", WS)
    add_title(slide, "The chat-protocol convention - type + envelope", top=Inches(1.15))
    add_subtitle(slide, "Most production WS apps use a tiny JSON envelope with a 'type' field both sides agree on.",
                 top=Inches(1.95))
    panel_w = Inches(6.0); panel_h = Inches(4.3); top = Inches(2.5)
    # CLIENT -> SERVER
    lx = Inches(0.6)
    lcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, top, panel_w, panel_h)
    lcard.fill.solid(); lcard.fill.fore_color.rgb = BG_PANEL
    lcard.line.color.rgb = WS; lcard.line.width = Pt(1.0); lcard.shadow.inherit = False
    lh = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.15), panel_w - Inches(0.5), Inches(0.4))
    lhr = lh.text_frame.paragraphs[0].add_run(); lhr.text = "CLIENT  ->  SERVER"
    lhr.font.name = FONT_HEAD; lhr.font.size = Pt(13); lhr.font.bold = True; lhr.font.color.rgb = WS
    c2s = """{ "type": "join",   "order_id": 123 }

{ "type": "msg",    "text": "I'm in 5C, buzzer broken" }

{ "type": "typing", "is_typing": true }

{ "type": "ping" }                       # heartbeat
"""
    tx = slide.shapes.add_textbox(lx + Inches(0.25), top + Inches(0.65), panel_w - Inches(0.5), panel_h - Inches(0.9))
    tf = tx.text_frame
    for i, line in enumerate(c2s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_CODE; r.font.size = Pt(11); r.font.color.rgb = TEXT
    # SERVER -> CLIENT
    rx = Inches(6.8)
    rcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, top, panel_w, panel_h)
    rcard.fill.solid(); rcard.fill.fore_color.rgb = BG_PANEL
    rcard.line.color.rgb = SSE; rcard.line.width = Pt(1.0); rcard.shadow.inherit = False
    rh = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.15), panel_w - Inches(0.5), Inches(0.4))
    rhr = rh.text_frame.paragraphs[0].add_run(); rhr.text = "SERVER  ->  CLIENT"
    rhr.font.name = FONT_HEAD; rhr.font.size = Pt(13); rhr.font.bold = True; rhr.font.color.rgb = SSE
    s2c = """{ "type": "history",  "messages": [...] }     # on join

{ "type": "msg",      "from": "raj", "text": "...", "ts": 1700000000 }

{ "type": "typing",   "from": "raj", "is_typing": true }

{ "type": "presence", "user": "sam", "online": true }

{ "type": "pong" }                       # heartbeat reply
"""
    tx = slide.shapes.add_textbox(rx + Inches(0.25), top + Inches(0.65), panel_w - Inches(0.5), panel_h - Inches(0.9))
    tf = tx.text_frame
    for i, line in enumerate(s2c.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_CODE; r.font.size = Pt(10); r.font.color.rgb = TEXT
    add_text(slide,
             "Receiver dispatches on type. Same shape powers Slack, Discord, Linear, Figma comments, etc.",
             top=Inches(6.95), size=12, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ws_pubsub_cross_server():
    slide = add_blank_slide()
    add_pattern_badge(slide, "WEBSOCKET", WS)
    add_title(slide, "Past 1 process: pub-sub or nothing reaches across servers", top=Inches(1.15))
    add_subtitle(slide, "Sam connects to server-A; Raj to server-B. How does Raj's message reach Sam?",
                 top=Inches(1.95))
    # Diagram with native shapes
    # Three boxes: server-A, Redis, server-B, with Sam/Raj attached
    # Simpler: render flow
    diag_top = Inches(2.6)
    # Sam left
    sam = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), diag_top, Inches(2.0), Inches(0.7))
    sam.fill.solid(); sam.fill.fore_color.rgb = BG_PANEL; sam.line.color.rgb = WS; sam.line.width = Pt(1.0); sam.shadow.inherit = False
    sx = slide.shapes.add_textbox(Inches(0.8), diag_top + Inches(0.15), Inches(2.0), Inches(0.4))
    sp = sx.text_frame.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run(); sr.text = "Sam's phone"
    sr.font.name = FONT_HEAD; sr.font.size = Pt(13); sr.font.bold = True; sr.font.color.rgb = WS
    # Server A
    sa = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), diag_top, Inches(2.5), Inches(0.7))
    sa.fill.solid(); sa.fill.fore_color.rgb = BG_PANEL; sa.line.color.rgb = SSE; sa.line.width = Pt(1.0); sa.shadow.inherit = False
    sax = slide.shapes.add_textbox(Inches(3.5), diag_top + Inches(0.15), Inches(2.5), Inches(0.4))
    sap = sax.text_frame.paragraphs[0]; sap.alignment = PP_ALIGN.CENTER
    sar = sap.add_run(); sar.text = "WS server A"
    sar.font.name = FONT_HEAD; sar.font.size = Pt(13); sar.font.bold = True; sar.font.color.rgb = SSE
    # Redis center
    red = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(3.6), Inches(2.5), Inches(0.7))
    red.fill.solid(); red.fill.fore_color.rgb = BG_PANEL; red.line.color.rgb = WEBHOOK; red.line.width = Pt(1.5); red.shadow.inherit = False
    rdx = slide.shapes.add_textbox(Inches(6.4), Inches(3.75), Inches(2.5), Inches(0.4))
    rdp = rdx.text_frame.paragraphs[0]; rdp.alignment = PP_ALIGN.CENTER
    rdr = rdp.add_run(); rdr.text = "Redis pub-sub"
    rdr.font.name = FONT_HEAD; rdr.font.size = Pt(13); rdr.font.bold = True; rdr.font.color.rgb = WEBHOOK
    # Server B
    sb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), diag_top + Inches(2.0), Inches(2.5), Inches(0.7))
    sb.fill.solid(); sb.fill.fore_color.rgb = BG_PANEL; sb.line.color.rgb = SSE; sb.line.width = Pt(1.0); sb.shadow.inherit = False
    sbx = slide.shapes.add_textbox(Inches(9.3), diag_top + Inches(2.15), Inches(2.5), Inches(0.4))
    sbp = sbx.text_frame.paragraphs[0]; sbp.alignment = PP_ALIGN.CENTER
    sbr = sbp.add_run(); sbr.text = "WS server B"
    sbr.font.name = FONT_HEAD; sbr.font.size = Pt(13); sbr.font.bold = True; sbr.font.color.rgb = SSE
    # Raj far right
    raj = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.9), diag_top + Inches(2.0), Inches(1.3), Inches(0.7))
    raj.fill.solid(); raj.fill.fore_color.rgb = BG_PANEL; raj.line.color.rgb = WS; raj.line.width = Pt(1.0); raj.shadow.inherit = False
    rajx = slide.shapes.add_textbox(Inches(11.9), diag_top + Inches(2.15), Inches(1.3), Inches(0.4))
    rajp = rajx.text_frame.paragraphs[0]; rajp.alignment = PP_ALIGN.CENTER
    rajr = rajp.add_run(); rajr.text = "Raj"
    rajr.font.name = FONT_HEAD; rajr.font.size = Pt(13); rajr.font.bold = True; rajr.font.color.rgb = WS
    # Step list
    steps = (
        "1. Raj's message arrives at server B  (Raj's chat WS lives there)\n"
        "2. Server B publishes to Redis channel chat:order:123\n"
        "3. Server A is subscribed - it receives the published message\n"
        "4. Server A forwards it down Sam's open WS  ->  Sam sees the message"
    )
    add_text(slide, steps, top=Inches(5.5), size=14, color=TEXT)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ---- Decision matrix additions -----
def slide_cost_ranking():
    slide = add_blank_slide()
    add_pattern_badge(slide, "DECIDE", LLM)
    add_title(slide, "Cost ranking - cheapest to most expensive per client", top=Inches(1.15))
    add_subtitle(slide, "Runtime cost only. Engineering cost roughly tracks it too.",
                 top=Inches(1.95))
    rows = [
        ("1.", "Webhook receiver",  WEBHOOK,
         "Zero idle cost. Server only does work when events fire. Scales for free when nothing's happening."),
        ("2.", "Short polling",     POLLING,
         "N reqs/min/client. Each is cheap and stateless. Cost grows linearly with clients x rate."),
        ("3.", "Long polling",      POLLING,
         "One held connection per client, no CPU until data arrives. Memory similar to SSE."),
        ("4.", "SSE",               SSE,
         "One held connection per client. Push cost when events fire. Tiny heartbeat overhead."),
        ("5.", "WebSocket",         WS,
         "SSE cost plus framing buffers + send queues + app-level reconnect/heartbeat state + pub-sub fanout."),
    ]
    top = Inches(2.55); h = Inches(0.75); gap = Inches(0.1)
    for i, (rank, name, color, body) in enumerate(rows):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = color; card.line.width = Pt(1.0); card.shadow.inherit = False
        rx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.18), Inches(0.6), Inches(0.4))
        rr = rx.text_frame.paragraphs[0].add_run(); rr.text = rank
        rr.font.name = FONT_HEAD; rr.font.size = Pt(20); rr.font.bold = True; rr.font.color.rgb = color
        nx = slide.shapes.add_textbox(Inches(1.5), y + Inches(0.18), Inches(3.0), Inches(0.4))
        nr = nx.text_frame.paragraphs[0].add_run(); nr.text = name
        nr.font.name = FONT_HEAD; nr.font.size = Pt(16); nr.font.bold = True; nr.font.color.rgb = color
        bx = slide.shapes.add_textbox(Inches(4.7), y + Inches(0.18), Inches(7.7), Inches(0.5))
        bf = bx.text_frame; bf.word_wrap = True
        br = bf.paragraphs[0].add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(12); br.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_recipes_extended():
    slide = add_blank_slide()
    add_pattern_badge(slide, "DECIDE", LLM)
    add_title(slide, "'I'm building X' - 15 quick recipes", top=Inches(1.15))
    add_subtitle(slide, "Most common AI / web scenarios mapped to the right pattern.",
                 top=Inches(1.95))
    items = [
        ("ChatGPT-style typewriter response",            "SSE",       SSE),
        ("Slack-style team chat",                        "WebSocket", WS),
        ("Stripe payment confirmation",                  "Webhook",   WEBHOOK),
        ("Vercel-style live deploy logs",                "SSE",       SSE),
        ("Multiplayer browser game",                     "WebSocket", WS),
        ("Voice agent (audio in + audio out)",           "WebSocket", WS),
        ("GitHub PR opened -> trigger CI/agent",         "Webhook",   WEBHOOK),
        ("'Is my batch job done yet?'",                  "Polling",   POLLING),
        ("Real-time stock ticker",                       "SSE",       SSE),
        ("Google-Docs-style collaborative editing",      "WebSocket", WS),
        ("MCP server you're building",                   "SSE / Streamable HTTP", SSE),
        ("Push to a mobile app while it's CLOSED",       "APNs / FCM (not us)", DIM),
        ("File upload progress in one HTTP request",     "fetch streaming (not us)", DIM),
        ("Internal service-to-service event trigger",    "Webhook",   WEBHOOK),
        ("Figma-style live cursors",                     "WebSocket", WS),
    ]
    # Two columns
    top = Inches(2.55); col_w = Inches(6.05); row_h = Inches(0.4); gap = Inches(0.05)
    for i, (scen, pat, color) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (col_w + Inches(0.13)) * col
        y = top + (row_h + gap) * row
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        sx = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.05), col_w - Inches(2.0), row_h - Inches(0.1))
        sp = sx.text_frame.paragraphs[0]
        sr = sp.add_run(); sr.text = scen
        sr.font.name = FONT_BODY; sr.font.size = Pt(11); sr.font.color.rgb = TEXT
        px = slide.shapes.add_textbox(x + col_w - Inches(2.05), y + Inches(0.05), Inches(2.0) - Inches(0.15), row_h - Inches(0.1))
        pp = px.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
        pr = pp.add_run(); pr.text = pat
        pr.font.name = FONT_HEAD; pr.font.size = Pt(11); pr.font.bold = True; pr.font.color.rgb = color
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_scale_tiers():
    slide = add_blank_slide()
    add_pattern_badge(slide, "DECIDE", LLM)
    add_title(slide, "What changes at scale - by tier", top=Inches(1.15))
    add_subtitle(slide, "Same patterns. Different infrastructure decisions at each tier.",
                 top=Inches(1.95))
    rows = [
        ("Up to 100 concurrent",
         "Anything works. Single server. Pick by ergonomics."),
        ("100 - 10K",
         "Async I/O. ulimit -n 65535. Heartbeats. LB idle timeout 60-300s. Monitor connection counts."),
        ("10K - 100K",
         "Multiple processes. Pub-sub backbone (Redis / NATS / Kafka). Sticky LB sessions. Connection metrics & slow-consumer detection."),
        ("100K - 1M",
         "Specialised tools - Phoenix Channels, Centrifugo, Ably, Pusher, PartyKit. Kernel tuning. Separate gateway tier. Edge delivery."),
        ("Over 1M",
         "Past blog-post territory. Specialised infrastructure built by teams who think about this all day. Hire someone who's done it."),
    ]
    top = Inches(2.55); h = Inches(0.85); gap = Inches(0.1)
    for i, (tier, body) in enumerate(rows):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        tx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.15), Inches(3.5), Inches(0.4))
        tr = tx.text_frame.paragraphs[0].add_run(); tr.text = tier
        tr.font.name = FONT_HEAD; tr.font.size = Pt(15); tr.font.bold = True; tr.font.color.rgb = ACCENT
        bx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.5), Inches(11.6), Inches(0.4))
        bf = bx.text_frame; bf.word_wrap = True
        br = bf.paragraphs[0].add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(12); br.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ---- AI / MCP additions -----
def slide_ai_why_different():
    slide = add_blank_slide()
    add_pattern_badge(slide, "AI SCENARIO", LLM)
    add_title(slide, "Why AI apps push every pattern decision harder", top=Inches(1.15))
    add_subtitle(slide, "Three forces that make these choices matter more in AI than in CRUD apps.",
                 top=Inches(1.95))
    items = [
        ("LLM output is naturally a stream",
         "Static REST feels broken once you've seen ChatGPT type. Batch responses feel laggy even when they're objectively fast."),
        ("Tools can take seconds to minutes",
         "Agent runs web search (5s), reads 3 docs (15s), summarises (10s), writes report (30s). The user is staring at a screen the whole time."),
        ("Agents are increasingly triggered by external events",
         "GitHub fires -> agent reviews PR. Stripe fires -> agent processes refund. Slack fires -> agent answers. That's webhooks. Lots of webhooks."),
    ]
    top = Inches(2.7); h = Inches(1.2); gap = Inches(0.15)
    for i, (head, body) in enumerate(items):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = LLM; card.line.width = Pt(1.5); card.shadow.inherit = False
        nx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.18), Inches(0.7), Inches(0.5))
        np = nx.text_frame.paragraphs[0]
        nr = np.add_run(); nr.text = str(i + 1)
        nr.font.name = FONT_HEAD; nr.font.size = Pt(28); nr.font.bold = True; nr.font.color.rgb = LLM
        hx = slide.shapes.add_textbox(Inches(1.7), y + Inches(0.18), Inches(10.7), Inches(0.4))
        hr = hx.text_frame.paragraphs[0].add_run(); hr.text = head
        hr.font.name = FONT_HEAD; hr.font.size = Pt(15); hr.font.bold = True; hr.font.color.rgb = TEXT
        bx = slide.shapes.add_textbox(Inches(1.7), y + Inches(0.6), Inches(10.7), Inches(0.55))
        bf = bx.text_frame; bf.word_wrap = True
        br = bf.paragraphs[0].add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(13); br.font.color.rgb = MUTED
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_mcp_transports():
    slide = add_blank_slide()
    add_pattern_badge(slide, "AI SCENARIO", LLM)
    add_title(slide, "MCP transports - three flavours, all SSE-shaped at the top", top=Inches(1.15))
    add_subtitle(slide, "Model Context Protocol = how Claude Desktop / Cursor talk to your tool servers.",
                 top=Inches(1.95))
    rows = [
        ("stdio",
         "local",
         "Server is a local process. Messages flow over stdin/stdout. Used for filesystem, git, sqlite tools. Not really a network pattern."),
        ("HTTP + SSE  (older design)",
         "remote, 2 endpoints",
         "Client opens an SSE stream for server-to-client. Client sends via a separate HTTP POST. Two endpoints glued together."),
        ("Streamable HTTP  (current standard)",
         "remote, 1 endpoint",
         "Single endpoint. The response body is itself an SSE stream. Bidirectional via paired POSTs and SSE responses. Still SSE-shaped underneath."),
    ]
    top = Inches(2.55); h = Inches(1.15); gap = Inches(0.13)
    for i, (name, tag, body) in enumerate(rows):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = LLM; card.line.width = Pt(1.0); card.shadow.inherit = False
        nx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.15), Inches(8.0), Inches(0.4))
        nr = nx.text_frame.paragraphs[0].add_run(); nr.text = name
        nr.font.name = FONT_HEAD; nr.font.size = Pt(15); nr.font.bold = True; nr.font.color.rgb = LLM
        tgx = slide.shapes.add_textbox(Inches(8.8), y + Inches(0.18), Inches(3.8), Inches(0.3))
        tgp = tgx.text_frame.paragraphs[0]; tgp.alignment = PP_ALIGN.RIGHT
        tgr = tgp.add_run(); tgr.text = tag
        tgr.font.name = FONT_HEAD; tgr.font.size = Pt(11); tgr.font.italic = True; tgr.font.color.rgb = MUTED
        bx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.55), Inches(11.6), Inches(0.55))
        bf = bx.text_frame; bf.word_wrap = True
        br = bf.paragraphs[0].add_run(); br.text = body
        br.font.name = FONT_BODY; br.font.size = Pt(12); br.font.color.rgb = TEXT
    add_text(slide,
             "Why SSE? Tool calls take time. One request, many response events (progress + result). That's exactly SSE's shape.",
             top=Inches(6.5), size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ai_external_trigger():
    slide = add_blank_slide()
    add_pattern_badge(slide, "AI SCENARIO", LLM)
    add_title(slide, "AI scenario - external event triggers the agent", top=Inches(1.15))
    add_subtitle(slide, "Order delivered -> LiveOrder sends a personalised thank-you SMS via an LLM.",
                 top=Inches(1.95))
    # Flow boxes
    boxes = [
        ("Order service",  Inches(0.6),  Inches(3.2),  WEBHOOK),
        ("Post-delivery\nservice", Inches(3.0),  Inches(3.2),  SSE),
        ("Queue",          Inches(5.7),  Inches(3.2),  ACCENT),
        ("Agent\n(LLM call)", Inches(7.7),  Inches(3.2),  LLM),
        ("Twilio API",     Inches(10.4), Inches(3.2),  WEBHOOK),
    ]
    for label, x, y, color in boxes:
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(1.0))
        b.fill.solid(); b.fill.fore_color.rgb = BG_PANEL
        b.line.color.rgb = color; b.line.width = Pt(1.5); b.shadow.inherit = False
        tx = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(2.2), Inches(0.7))
        tf = tx.text_frame; tf.word_wrap = True
        for i, line in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = line
            r.font.name = FONT_HEAD; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = color
    # Edge labels under the boxes
    add_text(slide, "internal webhook  ->  verify+dedup+enqueue  ->  worker pulls  ->  LLM call (1-5s)  ->  send SMS",
             top=Inches(4.5), size=12, color=MUTED, align=PP_ALIGN.CENTER, italic_via=True)
    # Three places people get this wrong
    add_text(slide, "Three places people get this wrong:",
             top=Inches(5.1), size=14, color=ACCENT, bold=True)
    add_text(slide,
             "1.  Doing the LLM call inside the webhook handler. LLM takes 1-5s; sender times out at 5-10s; you get retries; you send duplicate SMS.\n"
             "2.  No dedup. Same delivered event processed twice = customer gets 2 SMSes and Twilio bills you twice.\n"
             "3.  No idempotency in the worker. Two workers race the same job. Wrap the LLM+send in a DB transaction that checks 'already sent?'.",
             top=Inches(5.45), size=12, color=TEXT)
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


def slide_ai_five_questions():
    slide = add_blank_slide()
    add_pattern_badge(slide, "AI SCENARIO", LLM)
    add_title(slide, "Five questions to ask when building an AI feature", top=Inches(1.15))
    add_subtitle(slide, "Walk through these in order. The answers pick the pattern for you.",
                 top=Inches(1.95))
    qs = [
        ("Where does the trigger come from?",
         "User action -> REST / SSE / WS    ·    External event -> WEBHOOK    ·    Time-based -> cron"),
        ("Is the output stream-friendly?",
         "Token text -> SSE    ·    Binary audio / video -> WEBSOCKET    ·    Single result <1s -> REST"),
        ("Can the user interrupt mid-response?",
         "Yes -> WEBSOCKET (need a client->server signal mid-stream)    ·    No -> SSE is enough"),
        ("How long does the work take?",
         "<2s REST    ·    2-30s SSE for progress    ·    30s-min: job + SSE + webhook on done    ·    Hours: webhook callback only"),
        ("How many concurrent users?",
         "<100 anything    ·    100-10K SSE sweet spot    ·    >10K pub-sub backbone, consider managed WS"),
    ]
    top = Inches(2.55); h = Inches(0.85); gap = Inches(0.1)
    for i, (q, a) in enumerate(qs):
        y = top + (h + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.13), h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5); card.shadow.inherit = False
        nx = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.15), Inches(0.5), Inches(0.4))
        nr = nx.text_frame.paragraphs[0].add_run(); nr.text = str(i + 1) + "."
        nr.font.name = FONT_HEAD; nr.font.size = Pt(18); nr.font.bold = True; nr.font.color.rgb = LLM
        qx = slide.shapes.add_textbox(Inches(1.45), y + Inches(0.15), Inches(11.0), Inches(0.35))
        qr = qx.text_frame.paragraphs[0].add_run(); qr.text = q
        qr.font.name = FONT_HEAD; qr.font.size = Pt(14); qr.font.bold = True; qr.font.color.rgb = ACCENT
        ax = slide.shapes.add_textbox(Inches(1.45), y + Inches(0.5), Inches(11.0), Inches(0.4))
        af = ax.text_frame; af.word_wrap = True
        ar = af.paragraphs[0].add_run(); ar.text = a
        ar.font.name = FONT_BODY; ar.font.size = Pt(12); ar.font.color.rgb = TEXT
    add_footer(slide, len(prs.slides.__iter__.__self__._sldIdLst))


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------
def main():
    slide_cover()
    slide_agenda()
    slide_cast()
    slide_liveorder_feature_map()
    slide_http_problem()
    slide_http_request_response_raw()
    slide_stateless_concept()
    slide_status_codes()
    slide_stateless_table()
    slide_http_versions()
    slide_four_patterns_intro()
    polling_section()
    webhook_section()
    sse_section()
    ws_section()
    decision_section()
    slide_ai_why_different()
    slide_ai_long_task()
    slide_ai_external_trigger()
    slide_ai_voice_agent()
    slide_mcp_transports()
    composition_section()
    slide_ai_five_questions()
    materials_section()
    slide_recap()
    slide_qa()

    prs.save(OUT)
    n_slides = len(prs.slides.__iter__.__self__._sldIdLst)
    size = OUT.stat().st_size
    print(f"wrote {OUT.relative_to(HERE.parent)}")
    print(f"  {n_slides} slides")
    print(f"  {size:,} bytes ({size / 1024:.0f} KB)")


if __name__ == "__main__":
    main()
